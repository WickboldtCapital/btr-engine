import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

def add_logo_to_slide(slide, left, top, width):
    try:
        if os.path.exists("Gemini_Generated_Image_.png"):
            slide.shapes.add_picture("Gemini_Generated_Image_.png", left, top, width=width)
    except Exception:
        pass

def create_pptx(ui, calc):
    borrower_company = ui.get('borrower_company', '')
    project_name = ui.get('project_name', '')
    project_address = ui.get('project_address', '')
    report_date = ui.get('report_date', '')
    
    units = calc.get('units', 1)
    total_arv = calc.get('total_arv', 0)
    total_project_basis = calc.get('total_project_basis', 0)
    developer_margin = calc.get('developer_margin', 0)
    cash_surplus = calc.get('cash_surplus', 0)
    total_gross_monthly_income = calc.get('total_gross_monthly_income', 0)
    annual_noi = calc.get('annual_noi', 0)
    total_monthly_pi = calc.get('total_monthly_pi', 0)
    monthly_cash_flow = calc.get('monthly_cash_flow', 0)
    monthly_cash_flow_per_door = calc.get('monthly_cash_flow_per_door', 0)
    actual_dscr = calc.get('actual_dscr', 0)
    target_dscr_rate = calc.get('target_dscr_rate', 1.2)
    total_land_default = calc.get('total_land_default', 0)
    total_hard_cost = calc.get('total_hard_cost', 0)
    blended_cost_per_sf = calc.get('blended_cost_per_sf', 0)
    total_indirect_costs = calc.get('total_indirect_costs', 0)
    total_vertical_soft = calc.get('total_vertical_soft', 0)
    btr_finance_closing = calc.get('btr_finance_closing', 0)

    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    add_logo_to_slide(slide, Inches(3.75), Inches(0.5), Inches(2.5))
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{borrower_company} | BTR Pro Forma Analysis"
    subtitle.text = f"Project: {project_name if project_name else 'TBD'}\nAddress: {project_address if project_address else 'TBD'}\nPrepared: {report_date}"

    # 2. Exec Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    add_logo_to_slide(slide, Inches(8.0), Inches(0.2), Inches(1.5))
    
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = f"Project Scale: {units} Unit(s)"
    
    p = tf.add_paragraph()
    p.text = f"Total Appraised Value (ARV): ${total_arv:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Total Project Basis: ${total_project_basis:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Built-in Developer Margin: ${developer_margin:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Net Cash Surplus at Refi: ${cash_surplus:,.0f}" if cash_surplus >= 0 else f"Retained Seed Capital: ${-cash_surplus:,.0f}"

    # 3. Operating Metrics
    slide = prs.slides.add_slide(bullet_slide_layout)
    add_logo_to_slide(slide, Inches(8.0), Inches(0.2), Inches(1.5))
    
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Operating & DSCR Metrics"
    tf = body_shape.text_frame
    tf.text = f"Gross Monthly Income: ${total_gross_monthly_income:,.0f}"
    
    p = tf.add_paragraph()
    p.text = f"Net Operating Income (NOI): ${annual_noi/12:,.0f} / mo"
    p = tf.add_paragraph()
    p.text = f"Permanent Debt Service (P&I): ${total_monthly_pi:,.0f} / mo"
    p = tf.add_paragraph()
    p.text = f"Monthly Net Cash Flow: ${monthly_cash_flow:,.0f} (${monthly_cash_flow_per_door:,.0f} per door)"
    p = tf.add_paragraph()
    p.text = f"Actual DSCR: {actual_dscr:.2f}x (Target: {target_dscr_rate:.2f}x)"

    # 4. Capital Ledger
    slide = prs.slides.add_slide(bullet_slide_layout)
    add_logo_to_slide(slide, Inches(8.0), Inches(0.2), Inches(1.5))
    
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Capital Outlay Breakdown"
    tf = body_shape.text_frame
    tf.text = f"Horizontal Land Basis: ${total_land_default:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Vertical Direct Hard Costs: ${total_hard_cost:,.0f} (${blended_cost_per_sf:.2f} / SF)"
    p = tf.add_paragraph()
    p.text = f"Indirects & GC Fee: ${total_indirect_costs:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Soft Costs & Utilities: ${total_vertical_soft:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Financing & Closing Costs: ${btr_finance_closing:,.0f}"

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        with open(tmp.name, "rb") as f:
            return f.read()

def create_slide_pdf(ui, calc):
    borrower_company = ui.get('borrower_company', '')
    project_name = ui.get('project_name', '')
    project_address = ui.get('project_address', '')
    report_date = ui.get('report_date', '')
    
    units = calc.get('units', 1)
    total_arv = calc.get('total_arv', 0)
    total_project_basis = calc.get('total_project_basis', 0)
    developer_margin = calc.get('developer_margin', 0)
    cash_surplus = calc.get('cash_surplus', 0)
    total_gross_monthly_income = calc.get('total_gross_monthly_income', 0)
    annual_noi = calc.get('annual_noi', 0)
    total_monthly_pi = calc.get('total_monthly_pi', 0)
    monthly_cash_flow = calc.get('monthly_cash_flow', 0)
    monthly_cash_flow_per_door = calc.get('monthly_cash_flow_per_door', 0)
    actual_dscr = calc.get('actual_dscr', 0)
    target_dscr_rate = calc.get('target_dscr_rate', 1.2)
    total_land_default = calc.get('total_land_default', 0)
    total_hard_cost = calc.get('total_hard_cost', 0)
    blended_cost_per_sf = calc.get('blended_cost_per_sf', 0)
    total_indirect_costs = calc.get('total_indirect_costs', 0)
    total_vertical_soft = calc.get('total_vertical_soft', 0)
    btr_finance_closing = calc.get('btr_finance_closing', 0)

    pdf = FPDF(orientation='L', unit='mm', format='Letter')
    
    def add_pdf_logo(x, y, w):
        try:
            if os.path.exists("Gemini_Generated_Image_.png"):
                pdf.image("Gemini_Generated_Image_.png", x=x, y=y, w=w)
        except Exception:
            pass

    # 1. Title Slide
    pdf.add_page()
    add_pdf_logo(x=105, y=15, w=70)
    pdf.set_font('Arial', 'B', 28)
    pdf.cell(0, 60, '', ln=1) 
    pdf.cell(0, 15, f"{borrower_company} | BTR Pro Forma Analysis", align='C', ln=1)
    pdf.set_font('Arial', 'I', 16)
    pdf.cell(0, 10, f"Project: {project_name if project_name else 'TBD'}", align='C', ln=1)
    pdf.cell(0, 10, f"Address: {project_address if project_address else 'TBD'}", align='C', ln=1)
    pdf.cell(0, 10, f"Prepared: {report_date}", align='C', ln=1)
    
    # 2. Exec Summary
    pdf.add_page()
    add_pdf_logo(x=230, y=5, w=40)
    pdf.set_font('Arial', 'B', 24)
    pdf.cell(0, 20, 'Executive Summary', ln=1)
    pdf.line(10, 30, 270, 30) 
    pdf.set_font('Arial', '', 18)
    pdf.ln(10)
    pdf.cell(0, 12, f"- Project Scale: {units} Unit(s)", ln=1)
    pdf.cell(0, 12, f"- Total Appraised Value (ARV): ${total_arv:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Total Project Basis: ${total_project_basis:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Built-in Developer Margin: ${developer_margin:,.0f}", ln=1)
    cash_text = f"- Net Cash Surplus at Refi: ${cash_surplus:,.0f}" if cash_surplus >= 0 else f"- Retained Seed Capital: ${-cash_surplus:,.0f}"
    pdf.cell(0, 12, cash_text, ln=1)
    
    # 3. Operating Metrics
    pdf.add_page()
    add_pdf_logo(x=230, y=5, w=40)
    pdf.set_font('Arial', 'B', 24)
    pdf.cell(0, 20, 'Operating & DSCR Metrics', ln=1)
    pdf.line(10, 30, 270, 30)
    pdf.set_font('Arial', '', 18)
    pdf.ln(10)
    pdf.cell(0, 12, f"- Gross Monthly Income: ${total_gross_monthly_income:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Net Operating Income (NOI): ${annual_noi/12:,.0f} / mo", ln=1)
    pdf.cell(0, 12, f"- Permanent Debt Service (P&I): ${total_monthly_pi:,.0f} / mo", ln=1)
    pdf.cell(0, 12, f"- Monthly Net Cash Flow: ${monthly_cash_flow:,.0f} (${monthly_cash_flow_per_door:,.0f} per door)", ln=1)
    pdf.cell(0, 12, f"- Actual DSCR: {actual_dscr:.2f}x (Target: {target_dscr_rate:.2f}x)", ln=1)
    
    # 4. Capital Ledger
    pdf.add_page()
    add_pdf_logo(x=230, y=5, w=40)
    pdf.set_font('Arial', 'B', 24)
    pdf.cell(0, 20, 'Capital Outlay Breakdown', ln=1)
    pdf.line(10, 30, 270, 30)
    pdf.set_font('Arial', '', 18)
    pdf.ln(10)
    pdf.cell(0, 12, f"- Horizontal Land Basis: ${total_land_default:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Vertical Direct Hard Costs: ${total_hard_cost:,.0f} (${blended_cost_per_sf:.2f} / SF)", ln=1)
    pdf.cell(0, 12, f"- Indirects & GC Fee: ${total_indirect_costs:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Soft Costs & Utilities: ${total_vertical_soft:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Financing & Closing Costs: ${btr_finance_closing:,.0f}", ln=1)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        pdf.output(tmp.name)
        with open(tmp.name, "rb") as f:
            return f.read()
