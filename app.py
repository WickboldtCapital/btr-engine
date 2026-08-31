import streamlit as st
import pandas as pd
from fpdf import FPDF
from pptx import Presentation
import tempfile
import requests
import os
import json
import math
from datetime import datetime

# Page Configuration
st.set_page_config(page_title="Wickboldt Capital | BTR Pro Forma Engine", layout="wide")

# ==========================================
# --- MASTER DEFAULTS DICTIONARY ---
# ==========================================
HARDCODED_DRIVERS = {
    "units": 1,
    "sqft": 1173,
    "aux_cost_mode": "Percentage of Heated Rate (%)",
    "aux_rate_pct": 50.0,
    "aux_fixed_cost_sf": 35.0,
    "structure_type": "Carport", 
    "struct_sqft": 213,
    "front_porch_sqft": 49,
    "back_porch_sqft": 0,
    "storage_sqft": 49,
    "additional_foundation_cost": 3000,
    
    "gross_monthly_rent": 1600,
    "target_dscr_rate": 1.20,
    "target_min_cashflow_per_door": 200.0,
    "vacancy_rate_pct": 5.0,
    "opex_rate_pct": 20.0,
    
    "const_ltv_pct": 80.0,
    "build_months": 7,
    "const_rate_pct": 7.50,
    
    # Construction Loan Closing Fees
    "const_closing_mode": "Flat Lump Sum",
    "const_closing_fee": 6000,
    "const_origination_fee": 3000.0,
    "const_appraisal_fee": 600.0,
    "const_title_fee": 1200.0,
    "const_survey_fee": 500.0,
    "const_legal_fee": 500.0,
    "const_misc_closing_fee": 200.0,
    
    # Construction Lender Stress Constraints
    "const_bank_rent": 1500,
    "const_bank_ltv_pct": 80.0,
    "const_bank_val_mode": "Gross Rent Multiplier (GRM)",
    "const_bank_grm": 10.0,
    "const_bank_opex_pct": 30.0,
    "const_bank_vac_pct": 5.0,
    "const_bank_dscr": 1.20,
    "const_bank_qual_rate_pct": 7.25,
    "const_bank_amort_yrs": 30,
    "reserve_accrual_pct": 100.0,
    
    "refi_ltv_pct": 80.0,
    "refi_term_years": 30,
    "base_refi_rate_pct": 7.00,
    "refi_closing_fee": 5000,
    "apply_buydown": True,
    "buydown_pts": 2.75,
    
    "gc_fee_mode": "Percentage of Hard Costs (%)",
    "gc_fee_pct": 7.0,
    "custom_gc_fee": 20000,
    
    # Horizontal Parametric Development
    "land_entry_mode": "Flat Lump Sum per Lot",
    "site_type_mode": "Auto-Calc LF (Based on Lot Frontage)",
    "manual_infra_lf": 0.0,
    "land_basis": 10000,
    "land_raw_cost": 50000,
    "frontage_per_lot": 50.0,
    "lot_depth": 120.0,
    "road_width": 24.0,
    "road_layout": "Double-Sided Street (Yields 2 Lots per LF)",
    "lot_fill_inches": 12,
    "road_fill_inches": 10,
    "fill_cost_cy": 18.0,
    "clearing_per_lot": 3000.0,
    "paving_type": "Concrete",
    "paving_cost_lf": 125.0,
    "drainage_type": "Curb & Subsurface Catch Basins",
    "drainage_cost_lf": 85.0,
    "water_main_lf": 45.0,
    "sewer_main_lf": 65.0,
    "electric_main_lf": 35.0,
    "gas_main_lf": 20.0,
    "detention_pond": 15000.0,
    "site_amenities": 4500.0,
    
    # Vertical Hookups / Tap Fees
    "water_tap_fee": 1500.0,
    "sewer_tie_in": 1500.0,
    "electric_drop": 1000.0,
    "gas_lateral": 0.0,
    "impact_fees": 500.0,
    
    "appraisal_mode": "Income Approach (DSCR Loan Sizing)",
    "target_grm": 10.0,
    
    # NAHB & Indirect Cost Breakdowns
    "cost_calc_mode": "Reverse-Engineer from Primary Comp",
    "base_direct_cost_sf": 75.0,
    "lot_cost_pct": 13.7,
    "profit_margin_pct": 11.0,
    "overhead_pct": 5.7,
    "sales_pct": 3.6,
    "finance_pct": 1.5,
    "indirect_permits_pct": 4.5,
    "indirect_temp_facilities_pct": 3.5,
    
    "comp_address": "",
    "comp_price": 195000,
    "comp_heated_sf": 1275,
    "comp_struct_sf": 213,
    "comp_front_sf": 49,
    "comp_back_sf": 49,
    "comp_storage_sf": 0
}

GLOBAL_DRIVERS = HARDCODED_DRIVERS.copy()
for key, value in GLOBAL_DRIVERS.items():
    if key not in st.session_state:
        st.session_state[key] = value

if "raw_api_data" not in st.session_state: st.session_state.raw_api_data = None


# ==========================================
# --- SIDEBAR (Live Underwriting Controls) ---
# ==========================================
st.sidebar.header("Master Model Drivers")

with st.sidebar.container():
    st.subheader("1. Project Scale")
    st.number_input("Number of Units (Doors)", min_value=1, step=1, format="%d", key="units")

with st.sidebar.container():
    st.subheader("2. Market Comp Valuation & Benchmark")
    st.radio("Comparable Data Entry Mode", ["Manual Entry", "RentCast Live API Fetch"], key="comp_entry_mode")

    if st.session_state.comp_entry_mode == "RentCast Live API Fetch":
        search_address = st.text_input("Property Address", placeholder="e.g. 1103 S Spruce St, Hammond, LA")
        with st.expander("⚙️ API Configuration (Optional)"):
            manual_key = st.text_input("RentCast API Key Override", type="password")

        if st.button("Fetch Live Data", use_container_width=True):
            if search_address:
                rentcast_key = manual_key or os.environ.get("RENTCAST_API_KEY") or st.secrets.get("RENTCAST_API_KEY", "")
                if rentcast_key:
                    try:
                        with st.spinner("Fetching live data from RentCast MLS..."):
                            api_url = "https://api.rentcast.io/v1/properties"
                            headers = {"X-Api-Key": rentcast_key.strip(), "accept": "application/json"}
                            response = requests.get(api_url, headers=headers, params={"address": search_address.strip()})
                            if response.status_code == 200:
                                data = response.json()
                                st.session_state.raw_api_data = data 
                                if data:
                                    prop = data[0]
                                    if prop.get("price") or prop.get("lastSalePrice"):
                                        st.session_state.comp_price = int(prop.get("price") or prop.get("lastSalePrice"))
                                    if prop.get("squareFootage"):
                                        st.session_state.comp_heated_sf = int(prop.get("squareFootage"))
                                    if prop.get("formattedAddress"):
                                        st.session_state.comp_address = str(prop.get("formattedAddress"))
                                    st.success(f"Imported: {st.session_state.comp_address}")
                                    st.rerun()
                                else:
                                    st.error("No exact match found.")
                    except Exception as e:
                        st.error(f"API Error: {e}")

    st.markdown("##### 2.1 Primary Comp Metrics")
    is_rentcast = st.session_state.comp_entry_mode == "RentCast Live API Fetch"
    st.text_input("Comparable Property Address", key="comp_address", disabled=is_rentcast)
    st.number_input("Comp Sale Price ($)", min_value=0, step=1000, key="comp_price", disabled=is_rentcast)
    st.number_input("Comp Heated SF", min_value=0, step=50, key="comp_heated_sf", disabled=is_rentcast)

    st.markdown("##### 2.2 Comp Auxiliary Spaces")
    c_aux1, c_aux2 = st.columns(2)
    c_aux1.number_input("Aux SF", min_value=0, step=25, format="%d", key="comp_struct_sf", help="Garage or Carport")
    c_aux2.number_input("Front Porch SF", min_value=0, step=10, format="%d", key="comp_front_sf")
    c_aux3, c_aux4 = st.columns(2)
    c_aux3.number_input("Back Porch SF", min_value=0, step=10, format="%d", key="comp_back_sf")
    c_aux4.number_input("Storage SF", min_value=0, step=5, format="%d", key="comp_storage_sf")

with st.sidebar.container():
    st.subheader("3. Physical Footprint & Aux Costs")
    st.number_input("Heated SqFt per Unit", min_value=0, step=50, format="%d", key="sqft")
    
    st.radio("Auxiliary Rate Method", ["Percentage of Heated Rate (%)", "Fixed Cost ($ / SF)"], key="aux_cost_mode")
    
    if st.session_state.aux_cost_mode == "Percentage of Heated Rate (%)":
        st.slider("Auxiliary Rate (% of Heated Cost)", min_value=10.0, max_value=100.0, step=5.0, key="aux_rate_pct", help="Values unheated carport/porch space as a direct percentage of the heated shell rate per square foot for both the Comp and Project.")
    else:
        st.slider("Auxiliary Fixed Cost ($ / SF)", min_value=15.0, max_value=90.0, step=1.0, key="aux_fixed_cost_sf")
    
    st.selectbox("Aux Structure Type", ["Carport", "Garage"], key="structure_type")
    st.number_input(f"{st.session_state.structure_type} SqFt per Unit", min_value=0, step=25, format="%d", key="struct_sqft")
    st.number_input("Front Porch SqFt", min_value=0, step=10, format="%d", key="front_porch_sqft")
    st.number_input("Back Porch SqFt", min_value=0, step=10, format="%d", key="back_porch_sqft")
    st.number_input("Storage Room SqFt", min_value=0, step=5, format="%d", key="storage_sqft")
        
    st.number_input("Additional Foundation / Elevation Cost ($)", min_value=0, step=500, format="%d", key="additional_foundation_cost")

with st.sidebar.container():
    st.subheader("4. Cost Target Mode (Reverse Engineer)")
    st.radio("Calculation Logic", ["Manual Set (Heated SF)", "Reverse-Engineer from Appraisal", "Reverse-Engineer from Primary Comp"], key="cost_calc_mode")
    if st.session_state.cost_calc_mode == "Manual Set (Heated SF)":
        st.slider("Direct Build Cost / SF ($)", min_value=40.0, max_value=150.0, step=1.0, key="base_direct_cost_sf")
    else:
        st.slider("Finished Lot Cost (%)", min_value=0.0, max_value=30.0, step=0.1, key="lot_cost_pct")
        st.slider("Builder Profit (Pre-tax) (%)", min_value=0.0, max_value=30.0, step=0.1, key="profit_margin_pct")
        st.slider("Overhead & General Expenses (%)", min_value=0.0, max_value=15.0, step=0.1, key="overhead_pct")
        st.slider("Sales & Marketing Commissions (%)", min_value=0.0, max_value=15.0, step=0.1, key="sales_pct")
        st.slider("Financing Costs (%)", min_value=0.0, max_value=10.0, step=0.1, key="finance_pct")

with st.sidebar.container():
    st.subheader("5. Horizontal Land Development")
    st.radio("Land Basis Entry Mode", ["Flat Lump Sum per Lot", "Detailed Horizontal Infrastructure (LF Parametrics)"], key="land_entry_mode")
    
    if st.session_state.land_entry_mode == "Flat Lump Sum per Lot":
        st.number_input("Finished Lot Basis per Lot ($)", min_value=0, step=1000, format="%d", key="land_basis")
    else:
        st.number_input("Raw Land Acquisition ($ Total)", min_value=0, step=5000, key="land_raw_cost")
        
        st.markdown("##### Subdivision Geometry & Sizing")
        c_geom1, c_geom2 = st.columns(2)
        c_geom1.number_input("Average Lot Frontage (ft)", min_value=10.0, step=5.0, key="frontage_per_lot")
        c_geom2.number_input("Average Lot Depth (ft)", min_value=50.0, step=10.0, key="lot_depth")
        
        st.radio("Infrastructure Scope", ["Auto-Calc LF (Based on Lot Frontage)", "Manual Total LF (Infill / Custom)"], key="site_type_mode")
        if st.session_state.site_type_mode == "Auto-Calc LF (Based on Lot Frontage)":
            st.radio("Street Layout / Orientation", ["Single-Sided Street (1 Lot per LF)", "Double-Sided Street (Yields 2 Lots per LF)"], key="road_layout")
        else:
            st.number_input("Total New Street & Main Ext. Length (LF)", min_value=0.0, step=10.0, key="manual_infra_lf", help="Set to 0 for infill lots where the road and mains already exist.")

        st.markdown("##### Earthwork & Fill Dirt Calculator")
        c_fill1, c_fill2 = st.columns(2)
        c_fill1.number_input("Lot Fill Depth (Inches)", min_value=0, max_value=48, step=6, key="lot_fill_inches", help="Assumes 6-inch compaction lifts.")
        c_fill2.number_input("Road Fill Depth (Inches)", min_value=0, max_value=48, step=5, key="road_fill_inches", help="Assumes 5-inch compaction lifts.")
        
        c_fill3, c_fill4 = st.columns(2)
        c_fill3.number_input("Fill Cost ($/CY Compacted)", min_value=0.0, step=1.0, key="fill_cost_cy")
        c_fill4.number_input("Base Clear/Grub ($/Lot)", min_value=0.0, step=500.0, key="clearing_per_lot")

        st.markdown("##### Roadway & Trenching Rates ($ / LF)")
        st.number_input("Road Width (Feet)", min_value=10.0, step=2.0, key="road_width")
        c_pave1, c_pave2 = st.columns(2)
        c_pave1.selectbox("Paving Material", ["Asphalt", "Concrete"], key="paving_type")
        c_pave2.number_input("Paving Rate ($/LF)", min_value=0.0, step=5.0, key="paving_cost_lf")
        
        c_drain1, c_drain2 = st.columns(2)
        c_drain1.selectbox("Street Drainage", ["Curb & Catch Basins", "Open Ditch / Swale"], key="drainage_type")
        c_drain2.number_input("Drainage Rate ($/LF)", min_value=0.0, step=5.0, key="drainage_cost_lf")
        
        c_ut1, c_ut2 = st.columns(2)
        c_ut1.number_input("Water Main Ext. ($/LF)", min_value=0.0, step=2.5, key="water_main_lf")
        c_ut2.number_input("Sewer Main Ext. ($/LF)", min_value=0.0, step=2.5, key="sewer_main_lf")
        
        c_ut3, c_ut4 = st.columns(2)
        c_ut3.number_input("Elec/Comm Trench ($/LF)", min_value=0.0, step=2.5, key="electric_main_lf")
        c_ut4.number_input("Gas Main Ext. ($/LF)", min_value=0.0, step=2.5, key="gas_main_lf")

        st.markdown("##### Site-Wide Amenities ($ Total)")
        st.number_input("Stormwater Detention Pond ($ Total)", min_value=0.0, step=2500.0, key="detention_pond")
        st.number_input("Street Lights, Signage & Landscaping ($ Total)", min_value=0.0, step=1000.0, key="site_amenities")

with st.sidebar.container():
    st.subheader("6. Soft Costs, Permitting & GC Fees")
    st.slider("Permits, Fees, & Insurance (% of Direct)", min_value=0.0, max_value=15.0, step=0.5, key="indirect_permits_pct")
    st.slider("Temporary Site Facilities (% of Direct)", min_value=0.0, max_value=15.0, step=0.5, key="indirect_temp_facilities_pct")
    
    st.markdown("##### On-Lot Utility Hookups & Tap Fees ($ per Door)")
    st.number_input("Water Meter & Tap Fee", min_value=0.0, step=100.0, key="water_tap_fee")
    st.number_input("Sewer Tie-in Fee", min_value=0.0, step=100.0, key="sewer_tie_in")
    st.number_input("Electric Meter & Service Drop", min_value=0.0, step=100.0, key="electric_drop")
    st.number_input("Gas Service Lateral", min_value=0.0, step=100.0, key="gas_lateral")
    st.number_input("Municipal Impact Fees", min_value=0.0, step=250.0, key="impact_fees")

    st.radio("GC Fee Structure", ["Percentage of Hard Costs (%)", "Consolidated Flat Fee ($ Total)"], key="gc_fee_mode")
    if st.session_state.gc_fee_mode == "Percentage of Hard Costs (%)":
        st.number_input("GC Management Fee (%)", min_value=0.0, max_value=50.0, step=0.5, key="gc_fee_pct")
    else:
        st.number_input("Total Consolidated GC Fee ($)", min_value=0, step=1000, format="%d", key="custom_gc_fee")

with st.sidebar.container():
    st.subheader("7. Construction Loan Terms")
    st.slider("Construction / Bank LTC (%)", min_value=60.0, max_value=100.0, step=5.0, key="const_ltv_pct")
    st.slider("Construction Duration (Months)", min_value=3, max_value=18, step=1, key="build_months")
    st.slider("Construction Loan Rate (%)", min_value=4.0, max_value=14.0, step=0.5, key="const_rate_pct")
    
    st.markdown("##### Closing Costs")
    st.radio("Closing Fee Entry Mode", ["Flat Lump Sum", "Detailed Enterprise Breakout"], key="const_closing_mode")
    if st.session_state.const_closing_mode == "Flat Lump Sum":
        st.number_input("Const Loan Closing Fee ($ total)", min_value=0, step=500, format="%d", key="const_closing_fee")
    else:
        with st.expander("💼 Detailed Closing Costs", expanded=True):
            st.number_input("Origination / Lender Fee ($)", min_value=0.0, step=250.0, key="const_origination_fee")
            st.number_input("Appraisal & Feasibility ($)", min_value=0.0, step=100.0, key="const_appraisal_fee")
            st.number_input("Title, Escrow & Insurance ($)", min_value=0.0, step=100.0, key="const_title_fee")
            st.number_input("Survey & Environmental ($)", min_value=0.0, step=100.0, key="const_survey_fee")
            st.number_input("Legal & Doc Prep ($)", min_value=0.0, step=100.0, key="const_legal_fee")
            st.number_input("Recording & Misc ($)", min_value=0.0, step=50.0, key="const_misc_closing_fee")

with st.sidebar.container():
    st.subheader("8. Const. Lender Limits")
    st.number_input("Bank Underwriting Rent ($)", min_value=0, step=50, format="%d", key="const_bank_rent")
    st.slider("Bank Underwriting LTV (%)", min_value=50.0, max_value=100.0, step=5.0, key="const_bank_ltv_pct")
    
    st.radio("Bank Valuation Method", ["Gross Rent Multiplier (GRM)", "DSCR Stress Test"], key="const_bank_val_mode")
    
    if st.session_state.const_bank_val_mode == "Gross Rent Multiplier (GRM)":
        st.number_input("Bank Underwriting GRM", min_value=4.0, max_value=25.0, value=float(GLOBAL_DRIVERS.get("const_bank_grm", 10.0)), step=0.1, key="const_bank_grm")
    else:
        st.slider("Bank Underwriting Vacancy (%)", min_value=0.0, max_value=15.0, step=1.0, key="const_bank_vac_pct")
        st.slider("Bank Underwriting OpEx (%)", min_value=15.0, max_value=50.0, step=1.0, key="const_bank_opex_pct")
        st.number_input("Bank Target DSCR", min_value=1.0, max_value=1.5, value=float(GLOBAL_DRIVERS.get("const_bank_dscr", 1.20)), step=0.05, key="const_bank_dscr")
        st.slider("Bank Qualifying Rate (%)", min_value=4.0, max_value=14.0, step=0.25, key="const_bank_qual_rate_pct")
        st.selectbox("Bank Amortization (Years)", [15, 20, 25, 30], key="const_bank_amort_yrs")

with st.sidebar.container():
    st.subheader("9. Operating Pro Forma (DSCR)")
    st.number_input("Gross Monthly Rental Income per Unit ($)", min_value=0, step=50, format="%d", key="gross_monthly_rent")
    st.number_input("Target Lender DSCR Rate", min_value=1.0, max_value=1.5, value=float(GLOBAL_DRIVERS.get("target_dscr_rate", 1.20)), step=0.05, key="target_dscr_rate")
    st.number_input("Min. Acceptable Cash Flow / Door ($)", min_value=0.0, max_value=1000.0, value=float(GLOBAL_DRIVERS.get("target_min_cashflow_per_door", 200.0)), step=25.0, key="target_min_cashflow_per_door")
    st.slider("Vacancy Rate (%)", min_value=0.0, max_value=15.0, step=1.0, key="vacancy_rate_pct")
    st.slider("Operating Expenses (OpEx) Rate of EGI (%)", min_value=15.0, max_value=50.0, step=1.0, key="opex_rate_pct")

with st.sidebar.container():
    st.subheader("10. Takeout Appraisal Methodology")
    st.radio("Valuation Mode", [
        "Sales Comp (Price/SF)", 
        "Income Approach (GRM)", 
        "Income Approach (DSCR Loan Sizing)",
        "Conservative (Lesser of GRM or DSCR)"
    ], key="appraisal_mode")
    
    if st.session_state.appraisal_mode in ["Income Approach (GRM)", "Conservative (Lesser of GRM or DSCR)"]:
        st.number_input("Gross Rent Multiplier (GRM)", min_value=4.0, max_value=25.0, value=float(GLOBAL_DRIVERS.get("target_grm", 10.0)), step=0.1, key="target_grm")
        
    if st.session_state.appraisal_mode == "Income Approach (DSCR Loan Sizing)":
        st.info("ARV is automatically derived from your Operating (NOI) and Refi Loan terms to exactly meet your Target DSCR and LTV constraints.")
    elif st.session_state.appraisal_mode == "Conservative (Lesser of GRM or DSCR)":
        st.info("ARV evaluates both the GRM and the DSCR max loan limit, automatically defaulting to whichever yields the lower valuation.")

with st.sidebar.container():
    st.subheader("11. Takeout Refinance Terms")
    st.slider("Refinance LTV (%)", min_value=60.0, max_value=85.0, step=5.0, key="refi_ltv_pct")
    st.selectbox("Amortization Term (Years)", [15, 20, 25, 30], key="refi_term_years")
    st.slider("Base Refi Interest Rate (%)", min_value=4.0, max_value=10.0, step=0.25, key="base_refi_rate_pct")
    st.number_input("Refinance Closing Fee ($ total)", min_value=0, step=250, format="%d", key="refi_closing_fee")
    st.checkbox("Apply Interest Rate Buydown Points?", key="apply_buydown")
    if st.session_state.apply_buydown:
        st.number_input("Discount Points", min_value=0.0, max_value=10.0, step=0.25, format="%.2f", key="buydown_pts")
        net_rate = max(0.01, (st.session_state.base_refi_rate_pct / 100.0) - (st.session_state.buydown_pts * 0.0025))
        st.markdown(f"📉 **Buydown Net Rate:** `{net_rate*100:.3f}%`")


# ==========================================
# --- SAFE VARIABLE EXTRACTION FOR MATH ---
# ==========================================
units = st.session_state.get("units", 1)
sqft = st.session_state.get("sqft", 1173)

# Core Auxiliary logic 
aux_cost_mode = st.session_state.get("aux_cost_mode", "Percentage of Heated Rate (%)")
aux_ratio = st.session_state.get("aux_rate_pct", 50.0) / 100.0
aux_fixed_cost_sf = st.session_state.get("aux_fixed_cost_sf", 35.0)

struct_sqft = st.session_state.get("struct_sqft", 213)
front_porch_sqft = st.session_state.get("front_porch_sqft", 49)
back_porch_sqft = st.session_state.get("back_porch_sqft", 0)
storage_sqft = st.session_state.get("storage_sqft", 49)
additional_foundation_cost = st.session_state.get("additional_foundation_cost", 3000)

gross_monthly_rent = st.session_state.get("gross_monthly_rent", 1600)
target_dscr_rate = st.session_state.get("target_dscr_rate", 1.20)
target_min_cashflow_per_door = st.session_state.get("target_min_cashflow_per_door", 200.0)
vacancy_rate = st.session_state.get("vacancy_rate_pct", 5.0) / 100.0
opex_rate = st.session_state.get("opex_rate_pct", 20.0) / 100.0
const_ltv = st.session_state.get("const_ltv_pct", 80.0) / 100.0
build_months = st.session_state.get("build_months", 7)
const_rate = st.session_state.get("const_rate_pct", 7.50) / 100.0

# Construction Closing Fee Logic
const_closing_mode = st.session_state.get("const_closing_mode", "Flat Lump Sum")
if const_closing_mode == "Detailed Enterprise Breakout":
    c_orig = st.session_state.get("const_origination_fee", 3000.0)
    c_appr = st.session_state.get("const_appraisal_fee", 600.0)
    c_title = st.session_state.get("const_title_fee", 1200.0)
    c_surv = st.session_state.get("const_survey_fee", 500.0)
    c_leg = st.session_state.get("const_legal_fee", 500.0)
    c_misc = st.session_state.get("const_misc_closing_fee", 200.0)
    const_closing_fee = c_orig + c_appr + c_title + c_surv + c_leg + c_misc
else:
    const_closing_fee = st.session_state.get("const_closing_fee", 6000)

const_bank_rent = st.session_state.get("const_bank_rent", 1500)
const_bank_ltv = st.session_state.get("const_bank_ltv_pct", 80.0) / 100.0
const_bank_val_mode = st.session_state.get("const_bank_val_mode", "Gross Rent Multiplier (GRM)")
const_bank_grm = st.session_state.get("const_bank_grm", 10.0)
const_bank_opex_pct = st.session_state.get("const_bank_opex_pct", 30.0) / 100.0
const_bank_vac_pct = st.session_state.get("const_bank_vac_pct", 5.0) / 100.0
const_bank_dscr = st.session_state.get("const_bank_dscr", 1.20)
const_bank_qual_rate = st.session_state.get("const_bank_qual_rate_pct", 7.25) / 100.0
const_bank_amort_yrs = st.session_state.get("const_bank_amort_yrs", 30)

refi_ltv = st.session_state.get("refi_ltv_pct", 80.0) / 100.0
refi_term_years = st.session_state.get("refi_term_years", 30)
base_refi_rate = st.session_state.get("base_refi_rate_pct", 7.00) / 100.0
refi_closing_fee = st.session_state.get("refi_closing_fee", 5000)
apply_buydown = st.session_state.get("apply_buydown", True)
buydown_pts = st.session_state.get("buydown_pts", 2.75)
net_refi_rate = max(0.01, base_refi_rate - (buydown_pts * 0.0025)) if apply_buydown else base_refi_rate

gc_fee_mode = st.session_state.get("gc_fee_mode", "Percentage of Hard Costs (%)")
gc_fee_pct = st.session_state.get("gc_fee_pct", 7.0) / 100.0
custom_gc_fee = st.session_state.get("custom_gc_fee", 20000)

# Soft Costs / Utility Tie-ins
water_tap_fee = st.session_state.get("water_tap_fee", 1500.0)
sewer_tie_in = st.session_state.get("sewer_tie_in", 1500.0)
electric_drop = st.session_state.get("electric_drop", 1000.0)
gas_lateral = st.session_state.get("gas_lateral", 0.0)
impact_fees = st.session_state.get("impact_fees", 500.0)

total_tie_in_per_door = water_tap_fee + sewer_tie_in + electric_drop + gas_lateral + impact_fees
total_vertical_soft = total_tie_in_per_door * units

# Horizontal Land Development Parametrics
land_entry_mode = st.session_state.get("land_entry_mode", "Flat Lump Sum per Lot")
if land_entry_mode == "Detailed Horizontal Infrastructure (LF Parametrics)":
    land_raw_cost = st.session_state.get("land_raw_cost", 50000)
    frontage_per_lot = st.session_state.get("frontage_per_lot", 50.0)
    lot_depth = st.session_state.get("lot_depth", 120.0)
    
    site_type_mode = st.session_state.get("site_type_mode", "Auto-Calc LF (Based on Lot Frontage)")
    if site_type_mode == "Auto-Calc LF (Based on Lot Frontage)":
        road_layout = st.session_state.get("road_layout", "Double-Sided Street (Yields 2 Lots per LF)")
        total_road_lf = (frontage_per_lot * units) if "Single" in road_layout else (frontage_per_lot * units) / 2.0
    else:
        total_road_lf = st.session_state.get("manual_infra_lf", 0.0)
        
    road_width = st.session_state.get("road_width", 24.0)
    paving_cost_lf = st.session_state.get("paving_cost_lf", 125.0)
    drainage_cost_lf = st.session_state.get("drainage_cost_lf", 85.0)
    water_main_lf = st.session_state.get("water_main_lf", 45.0)
    sewer_main_lf = st.session_state.get("sewer_main_lf", 65.0)
    electric_main_lf = st.session_state.get("electric_main_lf", 35.0)
    gas_main_lf = st.session_state.get("gas_main_lf", 20.0)
    
    # Volumetric Earthwork Calculation
    lot_fill_inches = st.session_state.get("lot_fill_inches", 12.0)
    road_fill_inches = st.session_state.get("road_fill_inches", 10.0)
    fill_cost_cy = st.session_state.get("fill_cost_cy", 18.0)
    clearing_per_lot = st.session_state.get("clearing_per_lot", 3000.0)
    
    lot_sqft_total = frontage_per_lot * lot_depth * units
    road_sqft_total = total_road_lf * road_width
    
    lot_fill_cy = (lot_sqft_total * (lot_fill_inches / 12.0)) / 27.0
    road_fill_cy = (road_sqft_total * (road_fill_inches / 12.0)) / 27.0
    total_fill_cy = lot_fill_cy + road_fill_cy
    
    fill_dirt_cost = total_fill_cy * fill_cost_cy
    base_clearing_cost = clearing_per_lot * units
    land_earthwork = base_clearing_cost + fill_dirt_cost
    
    land_paving = total_road_lf * paving_cost_lf
    land_drainage = total_road_lf * drainage_cost_lf
    land_water_main = total_road_lf * water_main_lf
    land_sewer_main = total_road_lf * sewer_main_lf
    land_electric_main = total_road_lf * electric_main_lf
    land_gas_main = total_road_lf * gas_main_lf
    
    detention_pond = st.session_state.get("detention_pond", 15000.0)
    site_amenities = st.session_state.get("site_amenities", 4500.0)
    
    total_land_detailed = (land_raw_cost + land_earthwork + land_paving + land_drainage + 
                           land_water_main + land_sewer_main + land_electric_main + land_gas_main + 
                           detention_pond + site_amenities)
                           
    land_basis = total_land_detailed / units if units > 0 else 0
    total_land_default = total_land_detailed
else:
    land_basis = st.session_state.get("land_basis", 15000)
    total_land_default = land_basis * units

appraisal_mode = st.session_state.get("appraisal_mode", "Income Approach (DSCR Loan Sizing)")
target_grm = st.session_state.get("target_grm", 10.0)

# NAHB & Indirect Cost Breakdowns
cost_calc_mode = st.session_state.get("cost_calc_mode", "Reverse-Engineer from Primary Comp")
base_direct_cost_sf_input = st.session_state.get("base_direct_cost_sf", 75.0)
lot_cost_pct = st.session_state.get("lot_cost_pct", 13.7) / 100.0
profit_margin_pct = st.session_state.get("profit_margin_pct", 11.0) / 100.0
overhead_pct = st.session_state.get("overhead_pct", 5.7) / 100.0
sales_pct = st.session_state.get("sales_pct", 3.6) / 100.0
finance_pct = st.session_state.get("finance_pct", 1.5) / 100.0

indirect_permits_pct = st.session_state.get("indirect_permits_pct", 4.5) / 100.0
indirect_temp_facilities_pct = st.session_state.get("indirect_temp_facilities_pct", 3.5) / 100.0

comp_price = st.session_state.get("comp_price", 195000)
comp_heated_sf = st.session_state.get("comp_heated_sf", 1275)
comp_struct_sf = st.session_state.get("comp_struct_sf", 213)
comp_front_sf = st.session_state.get("comp_front_sf", 49)
comp_back_sf = st.session_state.get("comp_back_sf", 49)
comp_storage_sf = st.session_state.get("comp_storage_sf", 0)


# ==========================================
# --- CORE PRE-RENDER MATHEMATICS ---
# ==========================================
aux_sqft_total = struct_sqft + front_porch_sqft + back_porch_sqft + storage_sqft
total_under_roof_sqft = sqft + aux_sqft_total

comp_aux_sqft = comp_struct_sf + comp_front_sf + comp_back_sf + comp_storage_sf
comp_total_sf = comp_heated_sf + comp_aux_sqft
raw_comp_price_sf = comp_price / comp_heated_sf if comp_heated_sf > 0 else 0

# The combined deduction percentage used to extract the NAHB Construction Budget
combined_deductions_pct = lot_cost_pct + profit_margin_pct + overhead_pct + sales_pct + finance_pct
target_const_budget_pct = max(0, 1.0 - combined_deductions_pct)

# Multiplier to extract Direct Hard Cost from Total Construction Budget
indirect_multiplier = 1.0 + indirect_permits_pct + indirect_temp_facilities_pct

# 1. Evaluate Primary Comp and Isolate Heated Shell Rate
if aux_cost_mode == "Percentage of Heated Rate (%)":
    effective_comp_heated_sf = comp_heated_sf + (comp_aux_sqft * aux_ratio)
    isolated_heated_rate = comp_price / effective_comp_heated_sf if effective_comp_heated_sf > 0 else 0
    comp_aux_rate_sf = isolated_heated_rate * aux_ratio
    comp_aux_value = comp_aux_sqft * comp_aux_rate_sf
else:
    comp_aux_rate_sf = aux_fixed_cost_sf
    comp_aux_value = comp_aux_sqft * comp_aux_rate_sf
    comp_isolated_heated_value = max(0, comp_price - comp_aux_value)
    isolated_heated_rate = comp_isolated_heated_value / comp_heated_sf if comp_heated_sf > 0 else 0

# 2. Derive Valuation & Hard Cost Target
if appraisal_mode in ["Income Approach (GRM)", "Income Approach (DSCR Loan Sizing)", "Conservative (Lesser of GRM or DSCR)"]:
    # A. Calculate GRM Value
    grm_arv = (gross_monthly_rent * 12) * target_grm
    
    # B. Calculate DSCR Sized Value
    unit_gross_annual = gross_monthly_rent * 12.0
    unit_noi = unit_gross_annual * (1.0 - vacancy_rate) * (1.0 - opex_rate)
    unit_max_annual_ds = unit_noi / target_dscr_rate
    
    refi_monthly_rate = net_refi_rate / 12.0
    refi_term_months = refi_term_years * 12
    if refi_monthly_rate > 0:
        max_unit_loan = (unit_max_annual_ds / 12.0) * ((1.0 - (1.0 + refi_monthly_rate)**-refi_term_months) / refi_monthly_rate)
    else:
        max_unit_loan = (unit_max_annual_ds / 12.0) * refi_term_months
        
    dscr_arv = max_unit_loan / refi_ltv if refi_ltv > 0 else 0
    
    # C. Select Final ARV based on User Mode
    if appraisal_mode == "Income Approach (GRM)":
        arv_per_unit = grm_arv
    elif appraisal_mode == "Income Approach (DSCR Loan Sizing)":
        arv_per_unit = dscr_arv
    else: # Conservative (Lesser of GRM or DSCR)
        arv_per_unit = min(grm_arv, dscr_arv)
else:
    arv_per_unit = 0 

if aux_cost_mode == "Percentage of Heated Rate (%)":
    effective_project_heated_sf = sqft + (aux_sqft_total * aux_ratio)
    
    if cost_calc_mode == "Manual Set (Heated SF)":
        direct_cost_sf = base_direct_cost_sf_input
        struct_cost_sf = direct_cost_sf * aux_ratio
        our_aux_cost_total = (aux_sqft_total * struct_cost_sf) + additional_foundation_cost
        target_heated_hard_cost = direct_cost_sf * sqft
        target_direct_hard_cost = target_heated_hard_cost + our_aux_cost_total
        comp_equivalent_arv = (isolated_heated_rate * sqft) + (aux_sqft_total * (isolated_heated_rate * aux_ratio)) + additional_foundation_cost
        if appraisal_mode == "Sales Comp (Price/SF)":
            arv_per_unit = comp_equivalent_arv
            
        total_construction_budget = target_direct_hard_cost * indirect_multiplier * (1.0 + gc_fee_pct) if gc_fee_mode == "Percentage of Hard Costs (%)" else (target_direct_hard_cost * indirect_multiplier) + custom_gc_fee
            
    elif cost_calc_mode == "Reverse-Engineer from Appraisal":
        total_construction_budget = arv_per_unit * target_const_budget_pct
        if gc_fee_mode == "Percentage of Hard Costs (%)":
            target_direct_hard_cost = total_construction_budget / (indirect_multiplier * (1.0 + gc_fee_pct))
        else:
            target_direct_hard_cost = (total_construction_budget - custom_gc_fee) / indirect_multiplier
            
        direct_cost_sf = max(0, (target_direct_hard_cost - additional_foundation_cost) / effective_project_heated_sf) if effective_project_heated_sf > 0 else 0
        struct_cost_sf = direct_cost_sf * aux_ratio
        our_aux_cost_total = (aux_sqft_total * struct_cost_sf) + additional_foundation_cost
        target_heated_hard_cost = direct_cost_sf * sqft
        comp_equivalent_arv = arv_per_unit
        
    else: # Reverse-Engineer from Primary Comp
        comp_equivalent_arv = (isolated_heated_rate * sqft) + (aux_sqft_total * (isolated_heated_rate * aux_ratio)) + additional_foundation_cost
        total_construction_budget = comp_equivalent_arv * target_const_budget_pct
        if gc_fee_mode == "Percentage of Hard Costs (%)":
            target_direct_hard_cost = total_construction_budget / (indirect_multiplier * (1.0 + gc_fee_pct))
        else:
            target_direct_hard_cost = (total_construction_budget - custom_gc_fee) / indirect_multiplier
            
        direct_cost_sf = max(0, (target_direct_hard_cost - additional_foundation_cost) / effective_project_heated_sf) if effective_project_heated_sf > 0 else 0
        struct_cost_sf = direct_cost_sf * aux_ratio
        our_aux_cost_total = (aux_sqft_total * struct_cost_sf) + additional_foundation_cost
        target_heated_hard_cost = direct_cost_sf * sqft
        if appraisal_mode == "Sales Comp (Price/SF)":
            arv_per_unit = comp_equivalent_arv

else: # Fixed Aux Cost Mode
    struct_cost_sf = aux_fixed_cost_sf
    our_aux_cost_total = (aux_sqft_total * struct_cost_sf) + additional_foundation_cost
    
    if appraisal_mode == "Sales Comp (Price/SF)":
        arv_per_unit = (isolated_heated_rate * sqft) + our_aux_cost_total

    if cost_calc_mode == "Manual Set (Heated SF)":
        direct_cost_sf = base_direct_cost_sf_input
        target_heated_hard_cost = direct_cost_sf * sqft
        target_direct_hard_cost = target_heated_hard_cost + our_aux_cost_total
        comp_equivalent_arv = 0 
        total_construction_budget = target_direct_hard_cost * indirect_multiplier * (1.0 + gc_fee_pct) if gc_fee_mode == "Percentage of Hard Costs (%)" else (target_direct_hard_cost * indirect_multiplier) + custom_gc_fee

    elif cost_calc_mode == "Reverse-Engineer from Appraisal":
        total_construction_budget = arv_per_unit * target_const_budget_pct
        if gc_fee_mode == "Percentage of Hard Costs (%)":
            target_direct_hard_cost = total_construction_budget / (indirect_multiplier * (1.0 + gc_fee_pct))
        else:
            target_direct_hard_cost = (total_construction_budget - custom_gc_fee) / indirect_multiplier
        target_heated_hard_cost = max(0, target_direct_hard_cost - our_aux_cost_total)
        direct_cost_sf = target_heated_hard_cost / sqft if sqft > 0 else 0
        comp_equivalent_arv = arv_per_unit
        
    else: # Reverse from Comp
        comp_equivalent_arv = (isolated_heated_rate * sqft) + our_aux_cost_total
        total_construction_budget = comp_equivalent_arv * target_const_budget_pct
        if gc_fee_mode == "Percentage of Hard Costs (%)":
            target_direct_hard_cost = total_construction_budget / (indirect_multiplier * (1.0 + gc_fee_pct))
        else:
            target_direct_hard_cost = (total_construction_budget - custom_gc_fee) / indirect_multiplier
        target_heated_hard_cost = max(0, target_direct_hard_cost - our_aux_cost_total)
        direct_cost_sf = target_heated_hard_cost / sqft if sqft > 0 else 0

# Unified component rates across UI display
front_porch_cost_sf = struct_cost_sf
back_porch_cost_sf = struct_cost_sf
storage_cost_sf = struct_cost_sf

struct_total_cost = struct_sqft * struct_cost_sf
front_porch_cost = front_porch_sqft * front_porch_cost_sf
back_porch_cost = back_porch_sqft * back_porch_cost_sf
storage_cost = storage_sqft * storage_cost_sf
outdoor_aux_subtotal_unit = front_porch_cost + back_porch_cost + storage_cost

heated_hard_cost = sqft * direct_cost_sf
blended_cost_per_sf = (heated_hard_cost + our_aux_cost_total) / total_under_roof_sqft if total_under_roof_sqft > 0 else 0

# Base project hard costs
total_hard_cost = target_direct_hard_cost * units
total_arv = arv_per_unit * units
loan_total = total_arv * refi_ltv

default_buydown_cost = loan_total * (buydown_pts / 100.0) if apply_buydown else 0

# Target Lot Benchmarks
reference_price = arv_per_unit if cost_calc_mode == 'Reverse-Engineer from Appraisal' else comp_equivalent_arv
target_lot_value_per_door = reference_price * lot_cost_pct
target_total_lot_value = target_lot_value_per_door * units
land_equity_captured = target_total_lot_value - total_land_default

# 5.2 Indirects Build Up
permits_insurance_cost = total_hard_cost * indirect_permits_pct
temp_facilities_cost = total_hard_cost * indirect_temp_facilities_pct
fee_basis = total_hard_cost + permits_insurance_cost + temp_facilities_cost
default_gc_fee = custom_gc_fee if gc_fee_mode == "Consolidated Flat Fee ($ Total)" else fee_basis * gc_fee_pct
total_indirect_costs = permits_insurance_cost + temp_facilities_cost + default_gc_fee


# =========================================================================
# --- ACCURATE CONSTRUCTION SIZING, DSCR CAP & CAPITALIZED INTEREST ---
# =========================================================================

cb_gross_annual_rent = const_bank_rent * units * 12.0

if const_bank_val_mode == "DSCR Stress Test":
    cb_noi = cb_gross_annual_rent * (1.0 - const_bank_vac_pct) * (1.0 - const_bank_opex_pct)
    cb_max_annual_ds = cb_noi / const_bank_dscr
    cb_monthly_rate = const_bank_qual_rate / 12.0
    cb_term_months = const_bank_amort_yrs * 12

    if cb_monthly_rate > 0:
        bank_stressed_value = (cb_max_annual_ds / 12.0) * ((1.0 - (1.0 + cb_monthly_rate)**-cb_term_months) / cb_monthly_rate)
    else:
        bank_stressed_value = (cb_max_annual_ds / 12.0) * cb_term_months
else:
    bank_stressed_value = cb_gross_annual_rent * const_bank_grm

actual_const_loan = bank_stressed_value * const_bank_ltv

total_const = total_hard_cost + total_indirect_costs
total_project_costs_ex_interest = total_land_default + total_const + total_vertical_soft + const_closing_fee

# --- NEW S-CURVE ITERATIVE SOLVER ---
if land_entry_mode == "Detailed Horizontal Infrastructure (LF Parametrics)":
    land_raw_total = st.session_state.get("land_raw_cost", 50000)
    infra_total = total_land_default - land_raw_total
else:
    land_raw_total = total_land_default
    infra_total = 0.0

total_draw_budget = infra_total + total_const + total_vertical_soft

# Initial Guess for loop
carry_int_base = actual_const_loan * 0.5 * const_rate * (build_months / 12.0)

# Iterate 5 times to perfectly converge the capitalized interest cost against drawn loan balance
for _ in range(5):
    total_construction_basis = total_project_costs_ex_interest + carry_int_base
    seed_capital = max(0, total_construction_basis - actual_const_loan)
    
    equity_remaining = seed_capital
    day_1_costs = land_raw_total + const_closing_fee
    
    if equity_remaining >= day_1_costs:
        equity_remaining -= day_1_costs
        drawn_loan_balance = 0.0
    else:
        drawn_loan_balance = day_1_costs - equity_remaining
        equity_remaining = 0.0
        
    actual_scurve_interest = 0.0
    cum_pct = 0.0
    d_schedule = []
    
    for m in range(1, build_months + 1):
        prev_pct = cum_pct
        cum_pct = math.pow(math.sin((math.pi / 2.0) * (m / build_months)), 2)
        month_pct = cum_pct - prev_pct
        month_draw = total_draw_budget * month_pct
        
        month_int = drawn_loan_balance * (const_rate / 12.0)
        actual_scurve_interest += month_int
        
        if equity_remaining >= month_draw:
            equity_remaining -= month_draw
            funded_by_loan = 0.0
        else:
            funded_by_loan = month_draw - equity_remaining
            equity_remaining = 0.0
            
        drawn_loan_balance += funded_by_loan + month_int 
        
        d_schedule.append({
            "Month": f"Month {m}",
            "Cum. % Complete": f"{cum_pct*100:.1f}%",
            "Monthly Draw ($)": month_draw,
            "Capitalized Interest ($)": month_int,
            "Total Drawn Balance ($)": drawn_loan_balance
        })
        
    carry_int_base = actual_scurve_interest

df_schedule = pd.DataFrame(d_schedule)

total_construction_basis = total_project_costs_ex_interest + carry_int_base
total_project_basis = total_construction_basis + refi_closing_fee + default_buydown_cost
seed_capital = max(0, total_construction_basis - actual_const_loan)


# --- REFINANCE & OPERATING UNDERWRITING ---
monthly_interest_rate = net_refi_rate / 12.0
total_payments = refi_term_years * 12
if monthly_interest_rate > 0:
    monthly_pi_per_unit = (loan_total / units) * (monthly_interest_rate * (1 + monthly_interest_rate)**total_payments) / ((1 + monthly_interest_rate)**total_payments - 1)
else:
    monthly_pi_per_unit = (loan_total / units) / total_payments if total_payments > 0 else 0

total_monthly_pi = monthly_pi_per_unit * units
total_gross_monthly_income = gross_monthly_rent * units
monthly_vacancy_loss = total_gross_monthly_income * vacancy_rate
annual_vacancy_loss = monthly_vacancy_loss * 12.0
effective_gross_monthly_income = total_gross_monthly_income - monthly_vacancy_loss
annual_egi = effective_gross_monthly_income * 12.0
annual_opex = annual_egi * opex_rate
annual_noi = annual_egi - annual_opex
monthly_noi = annual_noi / 12.0

annual_debt_service = total_monthly_pi * 12.0
actual_dscr = annual_noi / annual_debt_service if annual_debt_service > 0 else 0
dscr_variance = actual_dscr - target_dscr_rate
monthly_cash_flow = monthly_noi - total_monthly_pi
monthly_cash_flow_per_door = monthly_cash_flow / units if units > 0 else 0

net_cash_at_closing = loan_total - actual_const_loan - refi_closing_fee - default_buydown_cost
cash_surplus = net_cash_at_closing - seed_capital
retained_equity = total_arv - loan_total
day1_wealth = default_gc_fee + max(0.0, cash_surplus) + retained_equity

btr_finance_closing = carry_int_base + const_closing_fee + refi_closing_fee + default_buydown_cost
lot_benchmark = target_total_lot_value
developer_margin = total_arv - (lot_benchmark + total_hard_cost + total_indirect_costs + total_vertical_soft + btr_finance_closing)

# =========================================================================
# --- SCALING CALCULATIONS (1 VS 3 VS 6) ---
# =========================================================================
unit_land = land_basis
unit_hard = target_direct_hard_cost
unit_soft = (permits_insurance_cost + total_vertical_soft) / units
unit_fin = btr_finance_closing / units
unit_gc_baseline = (temp_facilities_cost + default_gc_fee) / units

# 1-Unit Baseline
cap_1 = unit_land + unit_hard + unit_soft + unit_fin + unit_gc_baseline
loan_1 = loan_total / units
surplus_1 = loan_1 - cap_1
eq_1 = (total_arv / units) - loan_1
wealth_1 = unit_gc_baseline + max(0, surplus_1) + eq_1

# 3-Unit Phase
opt_gc_3 = 20000
savings_per_door = unit_gc_baseline - (opt_gc_3 / 3)
cap_3 = (unit_land + unit_hard + unit_soft + unit_fin) * 3 + opt_gc_3
loan_3 = loan_1 * 3
surplus_3 = loan_3 - cap_3
eq_3 = eq_1 * 3
wealth_3 = opt_gc_3 + max(0, surplus_3) + eq_3

# 6-Unit Program
opt_gc_6 = 40000
cap_6 = (unit_land + unit_hard + unit_soft + unit_fin) * 6 + opt_gc_6
loan_6 = loan_1 * 6
surplus_6 = loan_6 - cap_6
eq_6 = eq_1 * 6
wealth_6 = opt_gc_6 + max(0, surplus_6) + eq_6

monthly_cf_6 = monthly_cash_flow_per_door * 6
annual_cf_6 = monthly_cf_6 * 12


def format_surplus(val):
    return f"+${val:,.0f}" if val >= 0 else f"-${-val:,.0f}"

# =========================================================================
# --- SENSITIVITY ANALYSIS / STRESS TESTING MATH ---
# =========================================================================
stress_scenarios = [
    {"Name": "1. Baseline Target (Pro Forma)", "ARV_Shift": 0.0, "Rate_Shift": 0.0, "Delay": 0},
    {"Name": "2. Mild Rate Hike (+0.50%)", "ARV_Shift": 0.0, "Rate_Shift": 0.005, "Delay": 0},
    {"Name": "3. Severe Rate Hike (+1.00%)", "ARV_Shift": 0.0, "Rate_Shift": 0.01, "Delay": 0},
    {"Name": "4. Mild Appraisal Miss (-5%)", "ARV_Shift": -0.05, "Rate_Shift": 0.0, "Delay": 0},
    {"Name": "5. Severe Appraisal Miss (-10%)", "ARV_Shift": -0.10, "Rate_Shift": 0.0, "Delay": 0},
    {"Name": "6. Supply Chain Delay (+2 Months)", "ARV_Shift": 0.0, "Rate_Shift": 0.0, "Delay": 2},
    {"Name": "7. The Perfect Storm (-10% ARV, +1% Rate, +2 Mos)", "ARV_Shift": -0.10, "Rate_Shift": 0.01, "Delay": 2},
]

stress_results = []

for scn in stress_scenarios:
    # Adjust ARV and Takeout Loan
    s_arv = total_arv * (1 + scn["ARV_Shift"])
    s_loan_total = s_arv * refi_ltv
    s_buydown = s_loan_total * (buydown_pts / 100.0) if apply_buydown else 0
    s_months = build_months + scn["Delay"]
    
    # Recalculate S-Curve Carry Interest
    s_carry = actual_const_loan * 0.5 * const_rate * (s_months / 12.0)
    for _ in range(3):
        s_basis_ex_refi = total_project_costs_ex_interest + s_carry
        s_seed = max(0, s_basis_ex_refi - actual_const_loan)
        
        eq_rem = s_seed
        if eq_rem >= day_1_costs:
            eq_rem -= day_1_costs
            d_bal = 0.0
        else:
            d_bal = day_1_costs - eq_rem
            eq_rem = 0.0
        
        s_int = 0.0
        cp = 0.0
        for m in range(1, s_months + 1):
            pp = cp
            cp = math.pow(math.sin((math.pi / 2.0) * (m / s_months)), 2)
            m_drw = total_draw_budget * (cp - pp)
            m_i = d_bal * (const_rate / 12.0)
            s_int += m_i
            
            if eq_rem >= m_drw:
                eq_rem -= m_drw
                f_loan = 0.0
            else:
                f_loan = m_drw - eq_rem
                eq_rem = 0.0
            d_bal += f_loan + m_i
        s_carry = s_int
        
    s_seed_final = max(0, total_project_costs_ex_interest + s_carry - actual_const_loan)
    s_net_cash = s_loan_total - actual_const_loan - refi_closing_fee - s_buydown
    s_surplus = s_net_cash - s_seed_final
    s_retained_eq = s_arv - s_loan_total
    s_wealth = default_gc_fee + max(0.0, s_surplus) + s_retained_eq
    
    # Recalculate DSCR
    s_net_rate = net_refi_rate + scn["Rate_Shift"]
    s_m_rate = s_net_rate / 12.0
    if s_m_rate > 0:
        s_pi = (s_loan_total / units) * (s_m_rate * (1 + s_m_rate)**total_payments) / ((1 + s_m_rate)**total_payments - 1)
    else:
        s_pi = (s_loan_total / units) / total_payments if total_payments > 0 else 0
        
    s_ds = s_pi * units * 12.0
    s_dscr = annual_noi / s_ds if s_ds > 0 else 0
    
    # Emoticon formatting
    dscr_str = f"🟢 {s_dscr:.2f}x" if s_dscr >= target_dscr_rate else f"🔴 {s_dscr:.2f}x"
    surplus_str = f"🟢 +${s_surplus:,.0f}" if s_surplus >= 0 else f"🔴 -${abs(s_surplus):,.0f}"
    
    stress_results.append({
        "Stress Scenario": scn["Name"],
        "Final ARV": f"${s_arv:,.0f}",
        "Refi Rate": f"{s_net_rate*100:.2f}%",
        "Build Time": f"{s_months} mos",
        "DSCR Status": dscr_str,
        "Net Cash Surplus / (Trapped)": surplus_str,
        "Total Day-1 Wealth": f"${s_wealth:,.0f}"
    })

df_stress = pd.DataFrame(stress_results)

# ==========================================
# --- PAGE HEADER ---
# ==========================================
st.title("🏗️ BTR Pro Forma Engine")
st.markdown("### Wickboldt Capital — *Today's Foundation. Tomorrow's Legacy.*")
st.divider()

# --- PROJECT INFO ---
st.markdown("### 📋 Project Information")
top_col1, top_col2, top_col3 = st.columns([2, 2, 1])
project_name = top_col1.text_input("Project Title", placeholder="e.g. Phase 1 - 24-Lot Build-to-Rent")
project_address = top_col2.text_input("Project Address", placeholder="e.g. Rogers Moore Parkway, Hammond, LA")
report_date = datetime.now().strftime("%B %d, %Y")

sub_col1, sub_col2, sub_col3 = st.columns([1, 1, 2])
project_beds = sub_col1.number_input("Beds per Unit", min_value=1, value=3, step=1)
project_baths = sub_col2.number_input("Baths per Unit", min_value=1.0, value=2.0, step=0.5)

# --- DYNAMIC EXECUTIVE SUMMARY ---
address_display = project_address if project_address else "[Project Address]"
capital_recovery_text = "nearly 100% capital recovery" if -5000 <= cash_surplus <= 5000 else (f"a complete capital recovery with a ${cash_surplus:,.0f} surplus" if cash_surplus > 5000 else f"a capital recovery requiring ${-cash_surplus:,.0f} in retained seed capital")

executive_summary_text = (
    f"**Executive Summary:**\n"
    f"This technical brief outlines the pro forma and operational mechanics for a {units}-unit Build-to-Rent (BTR) "
    f"single-family residential project located at {address_display}. Utilizing a {sqft:,.0f} SF heated footprint "
    f"with a {struct_sqft:,.0f} SF integrated {st.session_state.structure_type.lower()}, the model reverse-engineers "
    f"national production builder economics to establish a target appraisal value (ARV) of ${arv_per_unit:,.0f} per home. "
    f"By deploying Wickboldt Capital as the managing general contractor, the project successfully captures active "
    f"construction management revenue while executing a commercial takeout refinance.\n\n"
    f"With a {build_months}-month construction timeline factored into carrying costs, this structure achieves "
    f"{capital_recovery_text}, a compliant {actual_dscr:.2f}x DSCR generating ${monthly_cash_flow_per_door:,.0f}/month "
    f"in net passive cash flow per door, and scalable wealth creation totaling ${day1_wealth:,.0f} across the "
    f"{units}-unit build program."
)

st.info(executive_summary_text.replace("$", r"\$"))
st.divider()

ui_top_metrics = st.container()
ui_op_metrics = st.container()
ui_decision_dashboard = st.container()
st.divider()

# ==========================================
# --- 1. POPULATE EXECUTIVE METRICS DASHBOARD ---
# ==========================================
with ui_top_metrics:
    st.markdown("### 1. Project Capital & Valuation Metrics")
    col1, col2, col3, col4 = st.columns(4)
    
    # Metric 1: Takeout Loan vs Total Project Basis
    refi_vs_basis = loan_total - total_project_basis
    col1.metric("Takeout Loan Proceeds", f"${loan_total:,.0f}", f"{refi_vs_basis:+,.0f} vs Project Basis", delta_color="normal")
    
    # Metric 2: Day-1 Seed Capital Status
    if seed_capital > 0:
        col2.metric("Day-1 Seed Capital", f"${seed_capital:,.0f}", "-Requires Cash Reserves", delta_color="normal")
    else:
        col2.metric("Day-1 Seed Capital", "$0", "Fully Funded", delta_color="normal")
        
    # Metric 3 & 4
    if cash_surplus >= 0:
        deal_health_color = "normal" 
        surplus_title = "Tax-Free Cash Surplus"
        surplus_delta = "Capital Recovered"
    else:
        deal_health_color = "inverse" 
        surplus_title = "Trapped Seed Capital"
        surplus_delta = "Loss at Closing"
        
    col3.metric("Under-Roof Blended Cost", f"${blended_cost_per_sf:.2f} / SF", f"{total_under_roof_sqft:,} Total SF Under Roof")
    col4.metric(surplus_title, f"${cash_surplus:,.0f}", surplus_delta, delta_color=deal_health_color)

with ui_op_metrics:
    st.markdown("### 🏢 Operating & DSCR Metrics")
    op1, op2, op3, op4 = st.columns(4)
    
    cf_pass = monthly_cash_flow_per_door >= target_min_cashflow_per_door
    
    arv_label = "Derived Unit ARV (Price/SF)" if appraisal_mode == "Sales Comp (Price/SF)" else "Derived Unit ARV"
    op1.metric(arv_label, f"${arv_per_unit:,.0f}", f"${total_arv:,.0f} Total ARV")
    op2.metric("Actual DSCR Rate", f"{actual_dscr:.2f}x", f"Target: {target_dscr_rate:.2f}x", delta_color="normal" if actual_dscr >= target_dscr_rate else "inverse")
    op3.metric("Monthly Cash Flow", f"${monthly_cash_flow:,.0f} /mo", f"${monthly_cash_flow*12:,.0f} Annual", delta_color="normal" if cf_pass else "inverse")
    op4.metric("Monthly P&I Payment", f"${total_monthly_pi:,.0f} /mo", f"{refi_term_years}Yr @ {net_refi_rate*100:.3f}%")

with ui_decision_dashboard:
    st.markdown("### 🚦 Go/No-Go Investment Decision Dashboard")
    dscr_pass = actual_dscr >= target_dscr_rate
    cash_pass = cash_surplus >= 0
    
    dscr_light = "🟢 GREEN LIGHT" if dscr_pass else "🔴 RED LIGHT"
    cf_light = "🟢 GREEN LIGHT" if cf_pass else "🔴 RED LIGHT"
    cash_light = "🟢 GREEN LIGHT" if cash_pass else "🔴 RED LIGHT"

    dash_col1, dash_col2, dash_col3, dash_col4 = st.columns(4)
    dash_col1.metric("DSCR Underwriting Status", dscr_light, f"Actual: {actual_dscr:.2f}x (Target: {target_dscr_rate:.2f}x)", delta_color="normal" if dscr_pass else "inverse")
    dash_col2.metric("Cash Flow Status", cf_light, f"Actual: ${monthly_cash_flow_per_door:,.0f}/door (Target: ${target_min_cashflow_per_door:,.0f})", delta_color="normal" if cf_pass else "inverse")
    dash_col3.metric("Capital Recovery Status", cash_light, f"${cash_surplus:,.0f} at Refi Close", delta_color="normal" if cash_pass else "inverse")
    dash_col4.metric("Day-1 Wealth Creation", f"${day1_wealth:,.0f}", f"${day1_wealth/units:,.0f} per door", delta_color="normal")


# ==========================================
# --- 2. COMP AUDIT DASHBOARD ---
# ==========================================
st.markdown("### 2. Market Comp Valuation Audit")

comp_under_roof = comp_heated_sf + comp_aux_sqft
comp_blended_rate = comp_price / comp_under_roof if comp_under_roof > 0 else 0
comp_address_display = st.session_state.comp_address if st.session_state.comp_address else "Market Benchmark Comparable"

comp_summary_text = (
    f"**Comparable Valuation Analysis:**\n\n"
    f"The primary comparable (**{comp_address_display}**) sold for **${comp_price:,.0f}**, spanning **{comp_heated_sf:,.0f} SF** heated living area "
    f"and **{comp_aux_sqft:,.0f} SF** of auxiliary space (**{comp_under_roof:,.0f} SF** total under roof at **${comp_blended_rate:.2f}/SF** blended).\n\n"
    f"To prevent non-living areas from inflating shell budgets, the engine algebraically isolates space types. "
    f"This yields a true isolated heated shell rate of **${isolated_heated_rate:.2f}/SF** and an implied auxiliary rate of "
    f"**${comp_aux_rate_sf:.2f}/SF** (discounted at **{aux_ratio*100:.0f}%** of the shell rate), establishing a pristine baseline for project ARV modeling."
)
st.info(comp_summary_text.replace("$", r"\$"))

with st.expander("🧮 View Comp Math Audit & Raw API Data", expanded=False):
    
    st.markdown("**2.1 Raw Heated Rate (Total Price ÷ Heated SF Only)**")
    st.code(f"${comp_price:,.0f} ÷ {comp_heated_sf:,.0f} Heated SF = ${raw_comp_price_sf:.2f} / SF")

    st.markdown("**2.2 Raw Blended Rate (Total Price ÷ Total Under-Roof SF)**")
    st.code(f"${comp_price:,.0f} ÷ {comp_under_roof:,.0f} Total SF = ${comp_blended_rate:.2f} / SF (Blended)")
    
    if aux_cost_mode == "Percentage of Heated Rate (%)":
        st.markdown(f"**2.3 True Isolated Heated Shell Rate (Algebraic Extraction)**")
        st.caption(f"If we simply halved the blended rate for the aux space, the total price would fall short. To perfectly distribute the ${comp_price:,.0f} across the spaces at a 100% / {aux_ratio*100:.0f}% ratio, we solve for the 'Effective Equivalent' square footage:".replace("$", r"\$"))
        st.code(f"Step 1: Find Equivalent SF\n{comp_heated_sf:,.0f} Heated SF + ({comp_aux_sqft:,.0f} Aux SF × {aux_ratio:.2f}) = {effective_comp_heated_sf:,.1f} Eq. SF\n\nStep 2: Solve for True Heated Rate\n${comp_price:,.0f} ÷ {effective_comp_heated_sf:,.1f} Eq. SF = ${isolated_heated_rate:.2f} / Heated SF\n\nStep 3: Solve for True Aux Rate\n${isolated_heated_rate:.2f} × {aux_ratio:.2f} = ${comp_aux_rate_sf:.2f} / Aux SF")
    else:
        st.markdown(f"**2.3 True Isolated Heated Shell Rate (Fixed Aux Value: ${aux_fixed_cost_sf:.2f} / SF)**".replace("$", r"\$"))
        st.code(f"(${comp_price:,.0f} - ${comp_aux_value:,.0f} Aux Value) ÷ {comp_heated_sf:,.0f} SF = ${isolated_heated_rate:.2f} / SF")
        
    if st.session_state.raw_api_data and st.session_state.get("comp_entry_mode") == "RentCast Live API Fetch":
        st.json(st.session_state.raw_api_data)
st.divider()


# ==========================================
# --- 3. HORIZONTAL LAND DEVELOPMENT ---
# ==========================================
st.markdown("### 3. Horizontal Land Development & Lot Basis")

if land_equity_captured >= 0:
    land_equity_text = f"instantly capturing **${land_equity_captured:,.0f}** in raw land equity before vertical construction even begins."
else:
    land_equity_text = f"running **${-land_equity_captured:,.0f}** over the target benchmark. Vertical construction margins must absorb this difference."

horizontal_summary_text = (
    f"The reverse-engineered retail model allocates **{lot_cost_pct*100:.1f}%** of the ARV to the finished lot, establishing a "
    f"target benchmark value of **${target_lot_value_per_door:,.0f}** per door (**${target_total_lot_value:,.0f}** total for {units} units). "
    f"By self-developing the infrastructure, our actual horizontal out-of-pocket basis is **${land_basis:,.0f}** per door "
    f"(**${total_land_default:,.0f}** total), {land_equity_text}"
)
st.info(horizontal_summary_text.replace("$", r"\$"))

with st.expander("🚜 View Land Development & Infrastructure Budget", expanded=False):
    if land_entry_mode == "Detailed Horizontal Infrastructure (LF Parametrics)":
        horizontal_data = {
            "Cost Category": [
                "Raw Land Acquisition",
                f"Earthwork & Fill ({total_fill_cy:,.0f} CY @ ${fill_cost_cy:,.0f}/CY + Clearing)",
                f"Water Main Infrastructure ({total_road_lf:,.0f} LF @ ${water_main_lf:,.0f}/LF)",
                f"Sewer Main Infrastructure ({total_road_lf:,.0f} LF @ ${sewer_main_lf:,.0f}/LF)",
                f"Electrical/Comm Trenching ({total_road_lf:,.0f} LF @ ${electric_main_lf:,.0f}/LF)",
                f"Gas Main Infrastructure ({total_road_lf:,.0f} LF @ ${gas_main_lf:,.0f}/LF)",
                f"Roadways, Paving & Flatwork ({total_road_lf:,.0f} LF @ ${paving_cost_lf:,.0f}/LF)",
                f"Stormwater & Drainage ({total_road_lf:,.0f} LF @ ${drainage_cost_lf:,.0f}/LF + Pond)",
                "Site Amenities, Lights & Signage",
                "TOTAL ACTUAL HORIZONTAL BASIS",
                "TARGET FINISHED LOT VALUE (BENCHMARK)",
                "CAPTURED LAND EQUITY"
            ],
            "Total Project Cost": [
                f"${land_raw_cost:,.0f}",
                f"${land_earthwork:,.0f}",
                f"${land_water_main:,.0f}",
                f"${land_sewer_main:,.0f}",
                f"${land_electric_main:,.0f}",
                f"${land_gas_main:,.0f}",
                f"${land_paving:,.0f}",
                f"${land_drainage + detention_pond:,.0f}",
                f"${site_amenities:,.0f}",
                f"${total_land_default:,.0f}",
                f"${target_total_lot_value:,.0f}",
                f"${land_equity_captured:,.0f}"
            ],
            "Cost Per Door": [
                f"${land_raw_cost/units:,.0f}" if units > 0 else "$0",
                f"${land_earthwork/units:,.0f}" if units > 0 else "$0",
                f"${land_water_main/units:,.0f}" if units > 0 else "$0",
                f"${land_sewer_main/units:,.0f}" if units > 0 else "$0",
                f"${land_electric_main/units:,.0f}" if units > 0 else "$0",
                f"${land_gas_main/units:,.0f}" if units > 0 else "$0",
                f"${land_paving/units:,.0f}" if units > 0 else "$0",
                f"${(land_drainage + detention_pond)/units:,.0f}" if units > 0 else "$0",
                f"${site_amenities/units:,.0f}" if units > 0 else "$0",
                f"${land_basis:,.0f}",
                f"${target_lot_value_per_door:,.0f}",
                f"${land_equity_captured/units:,.0f}" if units > 0 else "$0"
            ]
        }
        st.dataframe(pd.DataFrame(horizontal_data), hide_index=True, use_container_width=True)
    else:
        horizontal_data = {
            "Cost Category": [
                "Total Flat Finished Lot Basis",
                "Target Finished Lot Value (Benchmark)",
                "Captured Land Equity"
            ],
            "Total Project Cost": [
                f"${total_land_default:,.0f}",
                f"${target_total_lot_value:,.0f}",
                f"${land_equity_captured:,.0f}"
            ],
            "Cost Per Door": [
                f"${land_basis:,.0f}",
                f"${target_lot_value_per_door:,.0f}",
                f"${land_equity_captured/units:,.0f}" if units > 0 else "$0"
            ]
        }
        st.dataframe(pd.DataFrame(horizontal_data), hide_index=True, use_container_width=True)

st.divider()


# ==========================================
# --- 4. GRANULAR DIRECT HARD COST BUILDUP & ROLL-UP ---
# ==========================================
st.markdown("### 4. Granular Vertical Direct Hard Cost Buildup")

granular_summary_text = (
    f"Below is the itemized cost buildup per trade division required to achieve the baseline "
    f"**${direct_cost_sf:.2f}/SF (${heated_hard_cost:,.0f})** heated living area cost and the "
    f"**${struct_cost_sf:.2f}/SF (${struct_total_cost:,.0f})** {st.session_state.structure_type.lower()} cost, "
    f"yielding a blended direct hard cost of **${blended_cost_per_sf:.2f}/SF (${target_direct_hard_cost:,.0f})** under roof per unit."
)
st.info(granular_summary_text.replace("$", r"\$"))

with st.expander("🔨 View Granular Trade Buildup & Master Roll-up", expanded=False):
    granular_mode = st.radio("Buildup Entry Mode", ["Auto-Proportional (Linked to Master Model)", "Manual Custom Entry (Bottom-Up)"], horizontal=True)

    # Exact NAHB 8 Stages & 36 Components Mapping
    raw_heated_divs = [
        ("I. SITE WORK", 7.6, True),
        ("- Building Permit Fees", 1.8, False),
        ("- Impact Fee", 1.5, False),
        ("- Water & Sewer Fees / Inspections", 1.5, False),
        ("- Architecture & Engineering", 1.5, False),
        ("- Site Work - Other", 1.3, False),

        ("II. FOUNDATIONS", 10.5, True),
        ("- Excavation, Foundation, Concrete, Backfill", 10.0, False),
        ("- Foundations - Other", 0.5, False),

        ("III. FRAMING", 16.6, True),
        ("- Framing (including roof)", 11.6, False),
        ("- Trusses", 3.0, False),
        ("- Sheathing", 1.5, False),
        ("- General Metal, Steel", 0.4, False),
        ("- Framing - Other", 0.1, False),

        ("IV. EXTERIOR FINISHES", 13.4, True),
        ("- Exterior Wall Finish", 5.7, False),
        ("- Roofing", 3.9, False),
        ("- Windows and Doors (incl. garage door)", 3.7, False),
        ("- Exterior Finishes - Other", 0.1, False),

        ("V. MAJOR SYSTEMS ROUGH-INS", 19.2, True),
        ("- Plumbing (except fixtures)", 6.3, False),
        ("- Electrical (except fixtures)", 6.4, False),
        ("- HVAC", 6.3, False),
        ("- Major Systems - Other", 0.2, False),

        ("VI. INTERIOR FINISHES", 24.1, True),
        ("- Insulation", 1.6, False),
        ("- Drywall", 3.3, False),
        ("- Interior Trims, Doors, and Mirrors", 3.0, False),
        ("- Painting", 2.6, False),
        ("- Lighting", 1.3, False),
        ("- Cabinets, Countertops", 4.4, False),
        ("- Appliances", 1.7, False),
        ("- Flooring", 3.6, False),
        ("- Plumbing Fixtures", 1.8, False),
        ("- Fireplace", 0.6, False),
        ("- Interior Finishes - Other", 0.2, False),

        ("VII. FINAL STEPS", 6.5, True),
        ("- Landscaping", 2.2, False),
        ("- Outdoor Structures (deck, patio, porches)", 1.1, False),
        ("- Driveway", 2.3, False),
        ("- Clean Up", 0.9, False),

        ("VIII. OTHER", 2.1, True),
        ("- Miscellaneous / Other", 2.1, False),
    ]

    raw_struct_divs = [("AUXILIARY: CARPORT / GARAGE", 35.00, True)]
    pdf_granular_data = []

    if granular_mode == "Auto-Proportional (Linked to Master Model)":
        st.markdown(f"#### 4.1 Heated Living Area ({sqft} SF @ ${direct_cost_sf:.2f} / SF)".replace("$", r"\$"))
        h_data = {"Division / Trade Level": [], "Live Cost / SF": [], "Per Unit Cost": []}
        for name, base_val, is_header in raw_heated_divs:
            live_sf = direct_cost_sf * (base_val / 100.0)
            if is_header:
                h_data["Division / Trade Level"].append(f"{name}")
                h_data["Live Cost / SF"].append(f"${live_sf:.2f}")
                h_data["Per Unit Cost"].append(f"${live_sf * sqft:,.0f}")
                pdf_granular_data.append((name, live_sf, live_sf * sqft, live_sf * sqft * units, True))
            else:
                h_data["Division / Trade Level"].append(f"   {name}")
                h_data["Live Cost / SF"].append(f"${live_sf:.2f}")
                h_data["Per Unit Cost"].append(f"${live_sf * sqft:,.0f}")
                pdf_granular_data.append((f"   {name}", live_sf, live_sf * sqft, live_sf * sqft * units, False))
        st.dataframe(pd.DataFrame(h_data), hide_index=True, use_container_width=True)
        
        st.markdown(f"#### 4.2 {st.session_state.structure_type} Auxiliary ({struct_sqft} SF @ ${struct_cost_sf:.2f} / SF)".replace("$", r"\$"))
        s_data = {"Component Level": [], "Live Cost / SF": [], "Per Unit Cost": []}
        for name, base_val, is_header in raw_struct_divs:
            live_sf = struct_cost_sf * (base_val / 35.0)
            s_data["Component Level"].append(name)
            s_data["Live Cost / SF"].append(f"${live_sf:.2f}")
            s_data["Per Unit Cost"].append(f"${live_sf * struct_sqft:,.0f}")
        st.dataframe(pd.DataFrame(s_data), hide_index=True, use_container_width=True)
    else:
        st.markdown(f"#### 4.1 Heated Living Area ({sqft} SF)")
        hl_heated = []
        current_header = ""
        for name, base, is_h in raw_heated_divs:
            if is_h:
                current_header = name
            else:
                hl_heated.append({"Category": current_header, "Division / Component": name.replace("- ", ""), "Cost / SF": direct_cost_sf * (base/100.0)})
        
        df_manual = pd.DataFrame(hl_heated)
        edited_h = st.data_editor(df_manual.drop(columns=["Category"]), column_config={"Division / Component": st.column_config.TextColumn(disabled=True), "Cost / SF": st.column_config.NumberColumn(format="$%.2f", min_value=0.0, step=0.5)}, hide_index=True, use_container_width=True)
        
        direct_cost_sf = edited_h["Cost / SF"].sum()
        df_manual["Cost / SF"] = edited_h["Cost / SF"]
        
        # Intelligently re-aggregate Manual Custom Items back into the 8 NAHB Headers for the PDF Engine
        for cat in df_manual["Category"].unique():
            cat_df = df_manual[df_manual["Category"] == cat]
            cat_sum = cat_df["Cost / SF"].sum()
            pdf_granular_data.append((cat, cat_sum, cat_sum * sqft, cat_sum * sqft * units, True))
            for _, row in cat_df.iterrows():
                live_sf = row["Cost / SF"]
                pdf_granular_data.append((f"   - {row['Division / Component']}", live_sf, live_sf * sqft, live_sf * sqft * units, False))
        
        st.markdown(f"#### 4.2 {st.session_state.structure_type} Auxiliary ({struct_sqft} SF)")
        hl_struct = [{"Component Level": name, "Cost / SF": struct_cost_sf * (base/35.0)} for name, base, is_h in raw_struct_divs]
        edited_s = st.data_editor(pd.DataFrame(hl_struct), column_config={"Component Level": st.column_config.TextColumn(disabled=True), "Cost / SF": st.column_config.NumberColumn(format="$%.2f", min_value=0.0, step=0.5)}, hide_index=True, use_container_width=True)
        struct_cost_sf = edited_s["Cost / SF"].sum()

    st.markdown("#### 4.3 Auxiliary, Foundation & Outdoor Living")
    p_data = {
        "Component": ["Front Porch", "Back Porch", "Storage Room", "Addit. Foundation / Elevation", "TOTAL AUX/OUTDOOR"],
        "Area (SF)": [f"{front_porch_sqft} SF", f"{back_porch_sqft} SF", f"{storage_sqft} SF", "Site Specific", "Total"],
        "Cost / SF": [f"${front_porch_cost_sf:.2f}", f"${back_porch_cost_sf:.2f}", f"${storage_cost_sf:.2f}", "-", "-"],
        "Total Amount": [f"${front_porch_cost:,.0f}", f"${back_porch_cost:,.0f}", f"${storage_cost:,.0f}", f"${additional_foundation_cost:,.0f}", f"${our_aux_cost_total:,.0f}"]
    }
    st.dataframe(pd.DataFrame(p_data), hide_index=True, use_container_width=True)

    # Master Direct Hard Cost Roll-up Table
    st.markdown("#### 4.4 Direct Hard Cost Master Roll-up")
    direct_rollup_df = pd.DataFrame({
        "Category / Major Sub-Assembly": [
            f"1. Heated Living Area ({sqft:,} SF)",
            f"2. {st.session_state.structure_type} Structure ({struct_sqft:,} SF)",
            f"3. Outdoor & Storage Spaces ({front_porch_sqft + back_porch_sqft + storage_sqft:,} SF)",
            "4. Site Foundation & Elevation Premium",
            "TOTAL DIRECT HARD COST BUDGET"
        ],
        "Effective Rate": [
            f"${direct_cost_sf:.2f} / Heated SF",
            f"${struct_cost_sf:.2f} / SF",
            f"${struct_cost_sf:.2f} / SF (Avg)",
            "Flat Lump Sum",
            f"${blended_cost_per_sf:.2f} / SF Under Roof"
        ],
        "Per Unit Cost ($)": [
            heated_hard_cost,
            struct_total_cost,
            outdoor_aux_subtotal_unit,
            additional_foundation_cost,
            target_direct_hard_cost
        ],
        "Project Total ($)": [
            heated_hard_cost * units,
            struct_total_cost * units,
            outdoor_aux_subtotal_unit * units,
            additional_foundation_cost * units,
            total_hard_cost
        ]
    })

    st.dataframe(
        direct_rollup_df.style.format({
            "Per Unit Cost ($)": "${:,.0f}",
            "Project Total ($)": "${:,.0f}"
        }),
        hide_index=True,
        use_container_width=True
    )
st.divider()


# ==========================================
# --- 5. BANK STRESS TEST BREAKDOWN ---
# ==========================================
st.markdown("### 5. Construction Lender Loan Cap & Stress Test")

if const_bank_val_mode == "DSCR Stress Test":
    pdf_dscr_info_text = (
        f"**Construction Loan Underwriting:** The bank determines the implied asset value based on a "
        f"**{const_bank_dscr:.2f}x DSCR stress test** (yielding **${bank_stressed_value:,.0f}**). "
        f"They apply their **{const_bank_ltv*100:.1f}% LTV** limit to offer a maximum loan of **${actual_const_loan:,.0f}**, "
        f"and subtract that from your Total Basis (**${total_construction_basis:,.0f}**) to determine your required Day-1 Seed Capital of **${seed_capital:,.0f}**."
    )
    st.info(pdf_dscr_info_text.replace("$", r"\$"))
    
    stress_test_data = {
        "Bank Underwriting Step": [
            "1. Gross Potential Rent (Bank Model)",
            "2. Less: Bank Vacancy & OpEx Deductions",
            "3. Bank Qualified Net Operating Income (NOI)",
            "4. Target Construction DSCR Constraint",
            "5. Bank Implied Asset Value (Value of the House)",
            f"6. Const. Lender Loan Value ({const_bank_ltv*100:.1f}% Bank LTV)",
            "7. Total Capitalized Construction Basis",
            "8. Required Seed Capital Reserve (Basis - Loan Value)"
        ],
        "Value": [
            f"${cb_gross_annual_rent:,.0f} / yr",
            f"-${cb_gross_annual_rent - cb_noi:,.0f} / yr",
            f"${cb_noi:,.0f} / yr",
            f"{const_bank_dscr:.2f}x",
            f"${bank_stressed_value:,.0f}",
            f"${actual_const_loan:,.0f}",
            f"${total_construction_basis:,.0f}",
            f"${seed_capital:,.0f}"
        ]
    }
else:
    pdf_grm_info_text = (
        f"**Construction Loan Underwriting:** The bank determines the implied asset value based on a "
        f"**{const_bank_grm:.1f}x Gross Rent Multiplier** (yielding **${bank_stressed_value:,.0f}**). "
        f"They apply their **{const_bank_ltv*100:.1f}% LTV** limit to offer a maximum loan of **${actual_const_loan:,.0f}**, "
        f"and subtract that from your Total Basis (**${total_construction_basis:,.0f}**) to determine your required Day-1 Seed Capital of **${seed_capital:,.0f}**."
    )
    st.info(pdf_grm_info_text.replace("$", r"\$"))
    
    stress_test_data = {
        "Bank Underwriting Step": [
            "1. Gross Potential Rent (Bank Model)",
            "2. Bank Underwriting GRM",
            "3. Bank Implied Asset Value (Value of the House)",
            f"4. Const. Lender Loan Value ({const_bank_ltv*100:.1f}% Bank LTV)",
            "5. Total Capitalized Construction Basis",
            "6. Required Seed Capital Reserve (Basis - Loan Value)"
        ],
        "Value": [
            f"${cb_gross_annual_rent:,.0f} / yr",
            f"{const_bank_grm:.1f}x",
            f"${bank_stressed_value:,.0f}",
            f"${actual_const_loan:,.0f}",
            f"${total_construction_basis:,.0f}",
            f"${seed_capital:,.0f}"
        ]
    }

with st.expander("🏦 View Bank Stress Test Metrics", expanded=False):
    st.dataframe(pd.DataFrame(stress_test_data), hide_index=True, use_container_width=True)
    if seed_capital > 0:
        warning_text = f"**Seed Capital Required:** The bank will fund **${actual_const_loan:,.0f}** based on the appraised value. Your Total Construction Basis is **${total_construction_basis:,.0f}**. You must bring **${seed_capital:,.0f}** in Day-1 Seed Capital (Reserves)."
        st.warning(warning_text.replace("$", r"\$"))
    else:
        success_text = f"**Fully Funded:** The bank's loan value of **${actual_const_loan:,.0f}** covers your entire construction basis. No Day-1 Seed Capital is required."
        st.success(success_text.replace("$", r"\$"))
st.divider()


# ==========================================
# --- 6. DEVELOPER CAPITAL & CONSTRUCTION LEDGER ---
# ==========================================
st.markdown("### 6. Developer Capital & Construction Ledger")

land_eq_text = (
    f"**Land Equity Capture:** The Pro Forma valuation benchmark values the lot at {lot_cost_pct*100:.1f}% "
    f"(**${lot_benchmark:,.0f}**), while the developer's actual out-of-pocket acquisition/development basis is "
    f"**${total_land_default:,.0f}**. This gap represents immediate land equity."
)
st.info(land_eq_text.replace("$", r"\$"))

with st.expander("💰 View Capital Ledgers & Transaction Scope", expanded=False):
    fin_details = ""
    if const_closing_mode == "Detailed Enterprise Breakout":
        fin_details = "Origination, Appraisal, Title, Survey, Legal"

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("#### Pro Forma Valuation Allocation")
        val_alloc_df = pd.DataFrame({
            "Component": [f"Finished Lot Benchmark ({lot_cost_pct*100:.1f}%)", "Adjusted BTR Hard Costs", "Indirects + GC Fee", "On-Lot Utilities & Soft Costs", "Finance, Closing & Buydown", "Built-in Developer Margin", "Total Appraised Value (ARV)"],
            "Amount": [f"${lot_benchmark:,.0f}", f"${total_hard_cost:,.0f}", f"${total_indirect_costs:,.0f}", f"${total_vertical_soft:,.0f}", f"${btr_finance_closing:,.0f}", f"${developer_margin:,.0f}", f"${total_arv:,.0f}"]
        })
        st.dataframe(val_alloc_df, hide_index=True, use_container_width=True)
        
    with col2:
        st.markdown("#### Financing Breakdown")
        fin_breakdown_df = pd.DataFrame({
            "Component": ["Const. Loan Closing Fees", "Carrying Interest", "Perm. Takeout Fees", f"Rate Buydown ({buydown_pts:.2f} pts)"],
            "Amount": [f"${const_closing_fee:,.0f}", f"${carry_int_base:,.0f}", f"${refi_closing_fee:,.0f}", f"${default_buydown_cost:,.0f}"],
            "Details": [fin_details, f"({build_months} months @ {const_rate*100:.2f}% on S-Curve schedule)", "", ""]
        })
        st.dataframe(fin_breakdown_df, hide_index=True, use_container_width=True)

    st.markdown("#### Capital Outlay & Transaction Scope")
    transaction_df = pd.DataFrame({
        "Phase / Project Cash Event": ["1. Horizontal Dev & Land", "2. Vertical Construction Cost", "3. Utility Taps & Soft Costs", "4. Construction Financing Costs", "5. Permanent Refinance Takeout", "Total Project Capital Basis"],
        "Transaction Scope": ["Out-of-pocket acquisition and civil infrastructure", "Baseline Hard Costs + Indirects + GC Fee", "Water/sewer/electric/gas tie-ins, permits, builder's risk", f"Lender closing fees (${const_closing_fee:,.0f}) & accrued interest (${carry_int_base:,.0f})", f"Commercial fees (${refi_closing_fee:,.0f}) + {buydown_pts:.2f} pt rate buydown (${default_buydown_cost:,.0f})", "All-in total cost to build, finance, close, and stabilize"],
        "Actual Capital Outlay": [f"-${total_land_default:,.0f}", f"-${total_const:,.0f}", f"-${total_vertical_soft:,.0f}", f"-${const_closing_fee + carry_int_base:,.0f}", f"-${refi_closing_fee + default_buydown_cost:,.0f}", f"-${total_project_basis:,.0f}"]
    })
    st.dataframe(transaction_df, hide_index=True, use_container_width=True)
st.divider()


# ==========================================
# --- 7. OPERATING PERFORMANCE & DSCR ---
# ==========================================
st.markdown("### 7. Operating Performance Summary")

operating_summary_text = (
    f"**Stabilized Yield Analysis:**\n\n"
    f"Upon stabilization, the {units}-unit portfolio is projected to generate **${total_gross_monthly_income:,.0f}** in gross monthly rent. "
    f"After applying a **{vacancy_rate*100:.1f}%** vacancy allowance and a **{opex_rate*100:.1f}%** operating expense ratio, "
    f"the property yields **${monthly_noi:,.0f}** in monthly Net Operating Income (NOI). "
    f"Against a permanent debt service of **${total_monthly_pi:,.0f}** (modeled at **{net_refi_rate*100:.3f}%**), "
    f"the asset operates at a **{actual_dscr:.2f}x DSCR**, producing **${monthly_cash_flow_per_door:,.0f}** "
    f"in net passive cash flow per door every month."
)
st.info(operating_summary_text.replace("$", r"\$"))

with st.expander("🏢 View Stabilized Operating Pro Forma (DSCR)", expanded=False):
    dscr_summary_data = {
        "Pro Forma Line Item": [
            "Gross Potential Rent (GPR)", f"(-) Vacancy Loss @ {vacancy_rate*100:.1f}%", "= Effective Gross Income (EGI)", 
            f"(-) Operating Expenses (OpEx) @ {opex_rate*100:.1f}%", "= Net Operating Income (NOI)", "(-) Total Debt Service (P&I)", "= Net Cash Flow",
            "---", "Actual DSCR Rate", "Target Lender DSCR", "DSCR Variance"
        ],
        "Monthly": [
            f"${total_gross_monthly_income:,.2f}", f"-${monthly_vacancy_loss:,.2f}", f"${annual_egi / 12:,.2f}", 
            f"-${annual_opex / 12:,.2f}", f"${monthly_noi:,.2f}", f"-${total_monthly_pi:,.2f}", f"${monthly_cash_flow:,.2f}",
            "---", f"{actual_dscr:.2f}x", f"{target_dscr_rate:.2f}x", f"{dscr_variance:+.2f}x"
        ],
        "Annual": [
            f"${total_gross_monthly_income * 12:,.2f}", f"-${annual_vacancy_loss:,.2f}", f"${annual_egi:,.2f}", 
            f"-${annual_opex / 12:,.2f}", f"${annual_noi:,.2f}", f"${annual_debt_service:,.2f}", f"${monthly_cash_flow * 12:,.2f}",
            "---", f"{actual_dscr:.2f}x", f"{target_dscr_rate:.2f}x", f"{dscr_variance:+.2f}x"
        ]
    }
    st.dataframe(pd.DataFrame(dscr_summary_data), hide_index=True, use_container_width=True)
st.divider()


# ==========================================
# --- 8. REFINANCE WATERFALL & WEALTH ---
# ==========================================
st.markdown("### 8. Refinance Cash Waterfall & Day-1 Wealth")

waterfall_summary_text = (
    f"**Refinance Cash Waterfall:** This explicit Cash-Out Waterfall shows exactly how the "
    f"**${loan_total:,.0f}** Permanent Loan pays off the **${actual_const_loan:,.0f}** Construction phase "
    f"and **${refi_closing_fee + default_buydown_cost:,.0f}** in closing/buydown costs at the title company. "
)

if cash_surplus >= 0:
    waterfall_summary_text += f"After recovering your initial **${seed_capital:,.0f}** seed capital, the transaction yields a true **${cash_surplus:,.0f}** tax-free cash surplus."
else:
    waterfall_summary_text += f"After applying the net proceeds against your initial **${seed_capital:,.0f}** seed capital, **${-cash_surplus:,.0f}** remains as long-term retained capital in the deal."

st.info(waterfall_summary_text.replace("$", r"\$"))

with st.expander("💸 View Cash-Out Waterfall & Wealth Creation Breakdown", expanded=False):
    waterfall_data = {
        "Refinance Closing Line Item": [
            "(+) Permanent Takeout Loan Proceeds",
            "(-) Construction Loan Payoff (Prin. & Int.)",
            "(-) Refinance Closing Fees",
            f"(-) Rate Buydown Points Cost ({buydown_pts:.2f} pts)",
            "= Net Cash Distributed to Sponsor at Closing",
            "(-) Initial Sponsor Seed Capital Invested",
            "= Final Tax-Free Cash Surplus / (Trapped Capital)"
        ],
        "Amount ($)": [
            f"${loan_total:,.0f}",
            f"-${actual_const_loan:,.0f}",
            f"-${refi_closing_fee:,.0f}",
            f"-${default_buydown_cost:,.0f}",
            f"${net_cash_at_closing:,.0f}",
            f"-${seed_capital:,.0f}",
            f"${cash_surplus:,.0f}"
        ]
    }
    st.dataframe(pd.DataFrame(waterfall_data), hide_index=True, use_container_width=True)

    wealth_data = {
        "Pocket Component": ["Pocket 1: Active GC Fee Revenue", "Pocket 2: Tax-Free Cash Surplus at Close", "Pocket 3: Retained Asset Equity", "TOTAL DAY-1 CREATED VALUE"],
        "Value ($)": [f"${default_gc_fee:,.0f}", f"${cash_surplus:,.0f}", f"${retained_equity:,.0f}", f"${day1_wealth:,.0f}"]
    }
    st.dataframe(pd.DataFrame(wealth_data), hide_index=True, use_container_width=True)
st.divider()


# ==========================================
# --- 9. STRATEGY COMPARISON: RETAIL VS BTR ---
# ==========================================
st.markdown("### 9. Strategy Comparison: Retail Sell vs. Build-to-Rent")

retail_sales_costs = total_arv * 0.08
retail_total_invested = total_land_default + total_const + total_vertical_soft + retail_sales_costs
retail_net_cash = total_arv - retail_total_invested

btr_cash_text = f"yielding **${cash_surplus:,.0f}** tax-free" if cash_surplus >= 0 else f"leaving **${-cash_surplus:,.0f}** as retained capital"

strategy_comparison_text = (
    f"This comparison models the one-time taxable cash event of the National Builder model (yielding **${retail_net_cash:,.0f}** in taxable cash) "
    f"versus the wealth-creation mechanics of the Build-to-Rent model ({btr_cash_text}, while producing **${monthly_cash_flow_per_door:,.0f}/mo** in net cash flow and capturing **${retained_equity:,.0f}** in retained equity) "
    f"on the exact same project ({units} units, {sqft:,.0f} SF heated with {struct_sqft:,.0f} SF {st.session_state.structure_type.lower()}, ${total_arv:,.0f} ARV)."
)
st.info(strategy_comparison_text.replace("$", r"\$"))

with st.expander("⚖️ View Financial Strategy Comparison Matrix", expanded=False):
    strategy_data = {
        "Financial Metric": [
            "Land Basis",
            "Direct Hard Costs",
            "Total Construction Cost (incl. GC & Indirects)",
            "On-Lot Utilities & Soft Costs",
            "Sales & Transaction Closing Costs",
            "Total Capital Invested",
            "Gross Revenue / Permanent Loan Proceeds",
            "Net Cash Event at Closing",
            "Asset Retained on Balance Sheet?",
            "Ongoing Monthly Cash Flow"
        ],
        "National Builder (Retail Sell)": [
            f"${total_land_default:,.0f}",
            f"${total_hard_cost:,.0f} (Baseline)",
            f"${total_const:,.0f}",
            f"${total_vertical_soft:,.0f}",
            f"${retail_sales_costs:,.0f} (~8% Realtor & concessions)",
            f"${retail_total_invested:,.0f}",
            f"${total_arv:,.0f} (Retail sales price)",
            f"+${retail_net_cash:,.0f} (Taxable ordinary income)",
            "No",
            "$0"
        ],
        "Build-to-Rent (Commercial DSCR Takeout)": [
            f"${total_land_default:,.0f}",
            f"${total_hard_cost:,.0f}",
            f"${total_const:,.0f}",
            f"${total_vertical_soft:,.0f}",
            f"${btr_finance_closing:,.0f} (Const. finance, takeout, buydown)",
            f"${total_project_basis:,.0f}",
            f"${loan_total:,.0f} ({refi_ltv*100:.0f}% LTV DSCR Loan)",
            f"${cash_surplus:,.0f} (Seed Capital Retained)" if cash_surplus < 0 else f"+${cash_surplus:,.0f} (Tax-Free Cash Out)",
            f"Yes (${total_arv:,.0f} Asset, ${retained_equity:,.0f} Equity)",
            f"+${monthly_cash_flow_per_door:,.0f} / month / door"
        ]
    }

    st.dataframe(pd.DataFrame(strategy_data), hide_index=True, use_container_width=True)

    strategy_footnote_text = (
        f"***Note:** The BTR model captures ${default_gc_fee:,.0f} in GC Fees during the build. "
        f"The {build_months}-month construction timeline incurs carrying costs such that the takeout distribution "
        f"yields a net cash event of ${cash_surplus:,.0f} at closing. Total Day-1 Created Value equals ${day1_wealth:,.0f}.*"
    )
    st.caption(strategy_footnote_text.replace("$", r"\$"))
st.divider()


# ==========================================
# --- 10. MULTI-UNIT PHASE & ANNUAL PROGRAM ANALYSIS ---
# ==========================================
st.markdown("### 10. Multi-Unit Phase & Annual Program Analysis")

scaling_intro = (
    f"Scaling production unlocks substantial economies of scale on site supervision and management overhead. "
    f"This section compares a 3-house simultaneous phase against a 6-house annual build program (executed in two rolling phases of 3). "
    f"By consolidating supervision across each cluster, combined management fees are optimized down to **$20,000 per 3-house phase** "
    f"(${20000/3:,.0f} per house, saving **${savings_per_door:,.0f}** per home compared to single-unit builds)."
)
st.info(scaling_intro.replace("$", r"\$"))

with st.expander("📈 View Multi-Unit Scaling & Portfolio Optimization", expanded=False):
    scaling_data = {
        "Portfolio Metric": [
            "Total Units Built in Year 1",
            f"Land Basis (${unit_land:,.0f} / lot)",
            f"Adjusted Hard Costs (${unit_hard:,.0f} / house)",
            "Optimized Management Fees (GenCond + GC Fee)",
            f"On-Lot Utilities & Soft Costs (${unit_soft:,.0f} / house)",
            "Financing, Closing & Buydown",
            "Total Program Capital Invested",
            f"Permanent Commercial Loan Proceeds ({refi_ltv*100:.0f}% LTV)",
            "Net Program Cash Surplus at Close",
            f"Total Retained Asset Equity ({100 - refi_ltv*100:.0f}% ARV)",
            "Total Day-1 Wealth Created (Program)"
        ],
        "Single-House Baseline": [
            "1 Unit",
            f"${unit_land:,.0f}",
            f"${unit_hard:,.0f}",
            f"${unit_gc_baseline:,.0f}",
            f"${unit_soft:,.0f}",
            f"${unit_fin:,.0f}",
            f"${cap_1:,.0f}",
            f"${loan_1:,.0f}",
            format_surplus(surplus_1),
            f"${eq_1:,.0f}",
            f"${wealth_1:,.0f}"
        ],
        "3-House Simultaneous Phase": [
            "3 Units",
            f"${unit_land * 3:,.0f}",
            f"${unit_hard * 3:,.0f}",
            f"${opt_gc_3:,.0f}",
            f"${unit_soft * 3:,.0f}",
            f"${unit_fin * 3:,.0f}",
            f"${cap_3:,.0f}",
            f"${loan_3:,.0f}",
            format_surplus(surplus_3),
            f"${eq_3:,.0f}",
            f"${wealth_3:,.0f}"
        ],
        "6-House Annual Build (Year 1)": [
            "6 Units",
            f"${unit_land * 6:,.0f}",
            f"${unit_hard * 6:,.0f}",
            f"${opt_gc_6:,.0f}",
            f"${unit_soft * 6:,.0f}",
            f"${unit_fin * 6:,.0f}",
            f"${cap_6:,.0f}",
            f"${loan_6:,.0f}",
            format_surplus(surplus_6),
            f"${eq_6:,.0f}",
            f"${wealth_6:,.0f}"
        ],
        "Program Optimization Notes": [
            "Scaled velocity across fiscal year.",
            "Concurrent lot acquisition.",
            "Direct materials & labor.",
            "$20k per 3-house cluster ($6,667/house).",
            "Utility taps, impact fees, builder's risk.",
            f"Const. loans, {build_months}-mo carry interest, takeout pts.",
            "Total capital required for annual program.",
            f"${loan_1:,.0f} takeout loan per door at {net_refi_rate*100:.2f}% ({actual_dscr:.2f}x DSCR).",
            "Tax-free cash distributed above full seed return.",
            "Unencumbered equity across portfolio doors.",
            "Combined active GC fees, surplus, & equity."
        ]
    }

    st.dataframe(pd.DataFrame(scaling_data), hide_index=True, use_container_width=True)

    trapped_text = f"zero net capital trapped" if surplus_6 >= 0 else f"requiring ${-surplus_6:,.0f} in net capital trapped"
    surplus_extract_text = f"pulls out ${surplus_6:,.0f} in tax-free cash surplus" if surplus_6 >= 0 else f"retains ${-surplus_6:,.0f} as trapped seed capital"

    scaling_takeaway = (
        f"**Annual Build Program Takeaway (6 Houses in Year 1)**\n\n"
        f"Executing a 6-house annual program in two 3-house concurrent phases maximizes both corporate revenue and portfolio scale. "
        f"Wickboldt Capital earns **${opt_gc_6:,.0f}** in active GC management fees (saving **${savings_per_door * 6:,.0f}** overall), "
        f"the holding entity {surplus_extract_text} at refinance close, and the portfolio holds "
        f"**${eq_6:,.0f}** in unencumbered equity while generating **${monthly_cf_6:,.0f}/month** (**${annual_cf_6:,.0f}/year**) "
        f"in total net passive cash flow across all 6 doors at {trapped_text}."
    )
    st.success(scaling_takeaway.replace("$", r"\$"))
st.divider()


# ==========================================
# --- 11. ENTERPRISE S-CURVE DRAW SCHEDULE ---
# ==========================================
st.markdown("### 11. Enterprise S-Curve Draw Schedule & Actual Carry Interest")

scurve_summary = (
    f"**Capital Deployment Schedule:** Institutional lenders require actual drawdown forecasting. Rather than drawing the construction loan in a flat, straight line, "
    f"the model maps the **${total_draw_budget:,.0f}** physical development budget across a trigonometric S-Curve. "
    f"This realistically models slower initial site work, accelerated vertical framing, and tapering final finishes across the **{build_months}-month** timeline. "
    f"By using the exact cumulative distribution to dynamically calculate accrued interest, your True S-Curve Interest is mathematically locked at **${carry_int_base:,.0f}**, "
    f"ensuring an institutional-grade basis validation."
)
st.info(scurve_summary.replace("$", r"\$"))

with st.expander("📉 View S-Curve Draw Chart & Capitalization Schedule", expanded=False):
    # Render Chart
    chart_df = df_schedule.copy()
    chart_df = chart_df.set_index("Month")
    st.bar_chart(chart_df[["Monthly Draw ($)"]], height=300)

    # Render Table
    st.dataframe(df_schedule.style.format({
        "Monthly Draw ($)": "${:,.0f}",
        "Capitalized Interest ($)": "${:,.0f}",
        "Total Drawn Balance ($)": "${:,.0f}"
    }), hide_index=True, use_container_width=True)

st.divider()


# ==========================================
# --- 12. SENSITIVITY ANALYSIS (STRESS TEST) ---
# ==========================================
st.markdown("### 12. Sensitivity Analysis (Stress Testing)")

stress_test_summary = (
    f"**Lender Risk Mitigation Matrix:** This matrix dynamically stress-tests the active pro forma against severe market volatility. "
    f"It calculates the exact impact of commercial rate hikes, appraisal misses, and supply chain delays to prove whether the project "
    f"maintains a compliant DSCR (**> {target_dscr_rate:.2f}x**) and positive capital recovery under extreme duress."
)
st.info(stress_test_summary.replace("$", r"\$"))

with st.expander("🌩️ View Lender Risk Mitigation Matrix", expanded=False):
    st.dataframe(df_stress, hide_index=True, use_container_width=True)

st.divider()


# ==========================================
# --- 13. REVERSE-ENGINEERING BREAKDOWN ---
# ==========================================
if cost_calc_mode in ["Reverse-Engineer from Appraisal", "Reverse-Engineer from Primary Comp"]:
    st.markdown("### 13. Retail Comp & Appraisal Reverse-Engineering Breakdown")
    
    reference_price = arv_per_unit if cost_calc_mode == 'Reverse-Engineer from Appraisal' else comp_equivalent_arv
    ref_price_sf = reference_price / sqft if sqft > 0 else 0
    
    rev_eng_summary = (
        f"**Reverse-Engineering the Target Budget:**\n\nBased on the comp's isolated rates, your specific "
        f"{sqft} SF project has an estimated ARV of **${reference_price:,.0f}** (approximately "
        f"${ref_price_sf:,.2f} per heated square foot). We are reverse-engineering *this* value to isolate "
        f"the pure direct hard costs from finished land, builder corporate margins, indirects, and soft costs. "
        f"This ensures the physical build budget perfectly aligns with proven market economics."
    )
    st.info(rev_eng_summary.replace("$", r"\$"))
    
    with st.expander("⚙️ View Reverse-Engineering Budget Breakdown", expanded=False):
        breakdown_data = {
            "Cost Category": [
                f"Project Target ARV (Derived from Comp)", 
                "(-) Finished Lot Cost", 
                "(-) Builder Profit (Pre-tax)", 
                "(-) Overhead & General Expenses",
                "(-) Sales & Marketing", 
                "(-) Financing Costs", 
                "= Total Construction Budget (Direct + Indirect)", 
                "(-) Indirects (Permits, Temp Site, GC Fee)",
                "= Direct Hard Cost Budget",
                "(-) Fixed Auxiliary & Elevation Costs", 
                "= Available Budget for Heated Shell"
            ],
            "Value ($)": [
                f"${reference_price:,.0f}", 
                f"-${reference_price * lot_cost_pct:,.0f}", 
                f"-${reference_price * profit_margin_pct:,.0f}", 
                f"-${reference_price * overhead_pct:,.0f}", 
                f"-${reference_price * sales_pct:,.0f}", 
                f"-${reference_price * finance_pct:,.0f}", 
                f"${total_construction_budget:,.0f}", 
                f"-${total_indirect_costs:,.0f}",
                f"${target_direct_hard_cost:,.0f}",
                f"-${our_aux_cost_total:,.0f}", 
                f"${target_heated_hard_cost:,.0f}"
            ],
            "Description": [
                "Projected value of your build using the Comp's isolated rates.",
                "Land acquisition, grading, and utility infrastructure.",
                "Net profit margin retained by the homebuilding company.",
                "Office staff, software, vehicles, insurance, and administrative tools.",
                "Sales commissions and marketing costs to close the deal.",
                "Short-term interest paid on the builder's construction loan.",
                "Total budget available for all physical construction (Direct + Indirect).",
                "Permits, site facilities, and GC management fees.",
                "Total budget exclusively for physical structure and materials.",
                "Locked budget required for outdoor footprint and foundation elevation.",
                f"Yields exactly ${direct_cost_sf:.2f} / SF across {sqft} Heated SF."
            ]
        }
        st.dataframe(pd.DataFrame(breakdown_data), hide_index=True, use_container_width=True)
    st.divider()


# ==========================================
# --- 14. PDF GENERATION ENGINE ---
# ==========================================
st.markdown("### 🖨️ 14. Export Enterprise PDF Report")
pdf_detail_mode = st.radio(
    "Construction Budget Detail Level for PDF:",
    ["Full 36-Component Breakdown", "8 Major NAHB Categories Only", "High-Level Roll-up Only"],
    horizontal=True
)

class EnterpriseReport(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 16)
        self.cell(0, 8, 'Wickboldt Capital | BTR Pro Forma Report', border=0, ln=1, align='C')
        self.set_font('Arial', 'I', 10)
        self.cell(0, 6, "Today's Foundation. Tomorrow's Legacy.", border=0, ln=1, align='C')
        self.line(10, 25, 200, 25)
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()} | Prepared by: Stephen Wickboldt Jr. - Wickboldt Capital', 0, 0, 'C')

def create_pdf(detail_mode):
    pdf = EnterpriseReport()
    pdf.add_page()
    
    pdf.set_font("Arial", 'B', 12)
    p_name = project_name if project_name else "Untitled Development"
    p_address = project_address if project_address else "TBD"
    pdf.cell(0, 6, f"Project: {p_name}", ln=1)
    pdf.set_font("Arial", '', 10)
    pdf.cell(0, 6, f"Address: {p_address}", ln=1)
    pdf.cell(0, 6, f"Scale: {units} Unit(s) | {project_beds} Beds / {project_baths} Baths | Under-Roof: {total_under_roof_sqft:,} SF/Unit", ln=1)
    pdf.cell(0, 6, f"Date: {report_date}", ln=1)
    pdf.ln(5)

    # 1. EXECUTIVE NARRATIVE SUMMARY
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 1. Executive Narrative Summary", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    clean_summary = executive_summary_text.replace("**Executive Summary:**\n", "").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_summary)
    pdf.ln(5)

    # 2. EXECUTIVE METRICS & WEALTH CREATION
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 2. Capital Metrics & Wealth Creation", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    pdf.cell(100, 7, "Total Project Basis:", 0, 0)
    pdf.cell(90, 7, f"${total_project_basis:,.0f}", 0, 1, 'R')
    pdf.cell(100, 7, "Total Appraised Value (ARV):", 0, 0)
    pdf.cell(90, 7, f"${total_arv:,.0f}", 0, 1, 'R')
    pdf.cell(100, 7, "Const. Loan Proceeds (Payoff):", 0, 0)
    pdf.cell(90, 7, f"${actual_const_loan:,.0f}", 0, 1, 'R')
    pdf.cell(100, 7, f"Takeout Loan Proceeds ({refi_ltv*100:.0f}% LTV):", 0, 0)
    pdf.cell(90, 7, f"${loan_total:,.0f}", 0, 1, 'R')
    
    surplus_label = "Tax-Free Cash Surplus (At Closing):" if cash_surplus >= 0 else "Trapped Seed Capital (At Closing):"
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(100, 7, surplus_label, 0, 0)
    pdf.cell(90, 7, f"${cash_surplus:,.0f}", 0, 1, 'R')
    
    pdf.cell(100, 7, "Net Monthly Cash Flow per Door:", 0, 0)
    pdf.cell(90, 7, f"${monthly_cash_flow_per_door:,.0f}", 0, 1, 'R')
    
    pdf.cell(100, 7, "Total Day-1 Wealth Created:", 0, 0)
    pdf.cell(90, 7, f"${day1_wealth:,.0f}", 0, 1, 'R')
    pdf.ln(5)

    # 3. HORIZONTAL LAND DEVELOPMENT
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 3. Horizontal Land Development & Lot Basis", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    clean_horiz_summary = horizontal_summary_text.replace("**", "").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_horiz_summary)
    pdf.ln(3)

    pdf.set_font("Arial", 'B', 9)
    pdf.cell(100, 6, "Cost Category", 1, 0, 'C')
    pdf.cell(45, 6, "Total Project Cost", 1, 0, 'C')
    pdf.cell(45, 6, "Cost Per Door", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 8)
    for i in range(len(horizontal_data["Cost Category"])):
        cat = str(horizontal_data["Cost Category"][i])
        tot = str(horizontal_data["Total Project Cost"][i])
        per = str(horizontal_data["Cost Per Door"][i])
        
        # Make the total/summary rows bold
        if "TOTAL" in cat or "TARGET" in cat or "CAPTURED" in cat:
            pdf.set_font("Arial", 'B', 8)
            pdf.set_fill_color(240, 240, 240)
            fill = True
        else:
            pdf.set_font("Arial", '', 8)
            fill = False
            
        pdf.cell(100, 6, cat[:50], 1, 0, 'L', fill=fill)
        pdf.cell(45, 6, tot, 1, 0, 'R', fill=fill)
        pdf.cell(45, 6, per, 1, 1, 'R', fill=fill)
    
    pdf.ln(5)

    # 4. CONSTRUCTION BUILDUP SELECTION
    if detail_mode == "High-Level Roll-up Only":
        pdf.set_font("Arial", 'B', 12)
        pdf.set_fill_color(220, 220, 220)
        pdf.cell(0, 8, " 4. Vertical Direct Hard Cost Master Roll-up", ln=1, fill=True)
        pdf.set_font("Arial", 'B', 9)
        
        pdf.set_font("Arial", '', 10)
        granular_summary_pdf = (
            f"Below is the itemized cost buildup per trade division required to achieve the baseline "
            f"${direct_cost_sf:.2f}/SF (${heated_hard_cost:,.0f}) heated living area cost and the "
            f"${struct_cost_sf:.2f}/SF (${struct_total_cost:,.0f}) {st.session_state.structure_type.lower()} cost, "
            f"yielding a blended hard cost of ${blended_cost_per_sf:.2f}/SF (${target_direct_hard_cost:,.0f}) under roof per unit."
        )
        pdf.multi_cell(0, 6, granular_summary_pdf.encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(3)
        
        pdf.cell(85, 6, "Category / Major Sub-Assembly", 1, 0, 'C')
        pdf.cell(40, 6, "Effective Rate", 1, 0, 'C')
        pdf.cell(35, 6, "Per Unit Cost", 1, 0, 'C')
        pdf.cell(30, 6, "Project Total", 1, 1, 'C')
        
        for i, row in direct_rollup_df.iterrows():
            if i == len(direct_rollup_df) - 1:
                pdf.set_font("Arial", 'B', 9)
                pdf.set_fill_color(240, 240, 240)
                fill = True
            else:
                pdf.set_font("Arial", '', 9)
                fill = False
            
            pdf.cell(85, 6, str(row["Category / Major Sub-Assembly"]), 1, 0, 'L', fill=fill)
            pdf.cell(40, 6, str(row["Effective Rate"]), 1, 0, 'C', fill=fill)
            pdf.cell(35, 6, f"${row['Per Unit Cost ($)']:,.0f}", 1, 0, 'R', fill=fill)
            pdf.cell(30, 6, f"${row['Project Total ($)']:,.0f}", 1, 1, 'R', fill=fill)
            
    else:
        pdf.set_font("Arial", 'B', 12)
        pdf.set_fill_color(220, 220, 220)
        pdf.cell(0, 8, " 4. Granular Vertical Direct Hard Cost Buildup", ln=1, fill=True)
        pdf.set_font("Arial", 'B', 9)
        
        pdf.set_font("Arial", '', 10)
        granular_summary_pdf = (
            f"Below is the itemized cost buildup per trade division required to achieve the baseline "
            f"${direct_cost_sf:.2f}/SF (${heated_hard_cost:,.0f}) heated living area cost and the "
            f"${struct_cost_sf:.2f}/SF (${struct_total_cost:,.0f}) {st.session_state.structure_type.lower()} cost, "
            f"yielding a blended hard cost of ${blended_cost_per_sf:.2f}/SF (${target_direct_hard_cost:,.0f}) under roof per unit."
        )
        pdf.multi_cell(0, 6, granular_summary_pdf.encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(3)
        
        pdf.cell(90, 6, "Division / Trade Level", 1, 0, 'C')
        pdf.cell(30, 6, "Live Cost / SF", 1, 0, 'C')
        pdf.cell(35, 6, f"Per Unit ({sqft} SF)", 1, 0, 'C')
        pdf.cell(35, 6, "Project Total", 1, 1, 'C')
        
        for name, live_sf, unit_cost, proj_cost, is_header in pdf_granular_data:
            if detail_mode == "8 Major NAHB Categories Only" and not is_header:
                continue
                
            if is_header:
                pdf.set_font("Arial", 'B', 9)
                pdf.set_fill_color(240, 240, 240)
                pdf.cell(90, 6, name, 1, 0, 'L', fill=True)
                pdf.cell(30, 6, f"${live_sf:.2f}", 1, 0, 'R', fill=True)
                pdf.cell(35, 6, f"${unit_cost:,.0f}", 1, 0, 'R', fill=True)
                pdf.cell(35, 6, f"${proj_cost:,.0f}", 1, 1, 'R', fill=True)
            else:
                pdf.set_font("Arial", '', 9)
                pdf.cell(90, 6, name, 1, 0, 'L')
                pdf.cell(30, 6, f"${live_sf:.2f}", 1, 0, 'R')
                pdf.cell(35, 6, f"${unit_cost:,.0f}", 1, 0, 'R')
                pdf.cell(35, 6, f"${proj_cost:,.0f}", 1, 1, 'R')
            
    pdf.ln(5)

    # 5. CONSTRUCTION LENDER LIMITS
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 5. Construction Bank Limits & Loan Cap", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    if const_bank_val_mode == "DSCR Stress Test":
        pdf_dscr_info_text = (
            f"Construction Loan Underwriting: The bank determines the implied asset value based on a "
            f"{const_bank_dscr:.2f}x DSCR stress test (yielding ${bank_stressed_value:,.0f}). "
            f"They apply their {const_bank_ltv*100:.1f}% LTV limit to offer a maximum loan of ${actual_const_loan:,.0f}, "
            f"and subtract that from your Total Basis (${total_construction_basis:,.0f}) to determine your required Day-1 Seed Capital of ${seed_capital:,.0f}."
        )
        pdf.multi_cell(0, 6, pdf_dscr_info_text.encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(3)
        
        pdf.cell(100, 7, "Bank Assumed Rent & Target DSCR:", 0, 0)
        pdf.cell(90, 7, f"${const_bank_rent:,.0f}/mo | {const_bank_dscr:.2f}x DSCR", 0, 1, 'R')
        pdf.cell(100, 7, "Bank Vacancy & OpEx Deductions:", 0, 0)
        pdf.cell(90, 7, f"{const_bank_vac_pct*100:.1f}% Vac | {const_bank_opex_pct*100:.1f}% OpEx", 0, 1, 'R')
        pdf.cell(100, 7, "Bank Implied Asset Value (Refi Stress Test):", 0, 0)
        pdf.cell(90, 7, f"${bank_stressed_value:,.0f}", 0, 1, 'R')
    else:
        pdf_grm_info_text = (
            f"Construction Loan Underwriting: The bank determines the implied asset value based on a "
            f"{const_bank_grm:.1f}x Gross Rent Multiplier (yielding ${bank_stressed_value:,.0f}). "
            f"They apply their {const_bank_ltv*100:.1f}% LTV limit to offer a maximum loan of ${actual_const_loan:,.0f}, "
            f"and subtract that from your Total Basis (${total_construction_basis:,.0f}) to determine your required Day-1 Seed Capital of ${seed_capital:,.0f}."
        )
        pdf.multi_cell(0, 6, pdf_grm_info_text.encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(3)
        
        pdf.cell(100, 7, "Bank Assumed Rent & Target GRM:", 0, 0)
        pdf.cell(90, 7, f"${const_bank_rent:,.0f}/mo | {const_bank_grm:.1f}x GRM", 0, 1, 'R')
        pdf.cell(100, 7, "Bank Implied Asset Value (GRM Stress Test):", 0, 0)
        pdf.cell(90, 7, f"${bank_stressed_value:,.0f}", 0, 1, 'R')
    
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(130, 7, f"Const. Lender Loan Value ({const_bank_ltv*100:.1f}% Bank LTV):", 0, 0)
    pdf.cell(60, 7, f"${actual_const_loan:,.0f}", 0, 1, 'R')
    
    pdf.cell(130, 7, "Total Capitalized Construction Basis:", 0, 0)
    pdf.cell(60, 7, f"${total_construction_basis:,.0f}", 0, 1, 'R')
    
    pdf.cell(130, 7, "Required Seed Capital Reserve (Basis - Loan Value):", 0, 0)
    pdf.cell(60, 7, f"${seed_capital:,.0f}", 0, 1, 'R')
    pdf.ln(5)

    # 6. DEVELOPER CAPITAL & CONSTRUCTION LEDGER
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 6. Developer Capital & Construction Ledger", fill=True, ln=1)
    
    pdf.ln(3)
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(95, 6, "Pro Forma Valuation Allocation", 0, 0, 'L')
    pdf.cell(95, 6, "Financing Breakdown", 0, 1, 'L')
    
    pdf.set_font("Arial", '', 9)
    # Row 1
    pdf.cell(65, 5, f"Finished Lot Benchmark ({lot_cost_pct*100:.0f}%):", 0, 0, 'L')
    pdf.cell(30, 5, f"${lot_benchmark:,.0f}", 0, 0, 'R')
    pdf.cell(65, 5, "Const. Loan Closing Fees:", 0, 0, 'L')
    pdf.cell(30, 5, f"${const_closing_fee:,.0f}", 0, 1, 'R')
    # Row 2
    pdf.cell(65, 5, "Adjusted BTR Hard Costs:", 0, 0, 'L')
    pdf.cell(30, 5, f"${total_hard_cost:,.0f}", 0, 0, 'R')
    pdf.cell(65, 5, "Carrying Interest:", 0, 0, 'L')
    pdf.cell(30, 5, f"${carry_int_base:,.0f}", 0, 1, 'R')
    # Row 3
    pdf.cell(65, 5, "Indirects + GC Fee:", 0, 0, 'L')
    pdf.cell(30, 5, f"${total_indirect_costs:,.0f}", 0, 0, 'R')
    pdf.cell(65, 5, "Perm. Takeout Fees:", 0, 0, 'L')
    pdf.cell(30, 5, f"${refi_closing_fee:,.0f}", 0, 1, 'R')
    # Row 4
    pdf.cell(65, 5, "On-Lot Utilities & Soft Costs:", 0, 0, 'L')
    pdf.cell(30, 5, f"${total_vertical_soft:,.0f}", 0, 0, 'R')
    pdf.cell(65, 5, f"Rate Buydown ({buydown_pts:.2f} pts):", 0, 0, 'L')
    pdf.cell(30, 5, f"${default_buydown_cost:,.0f}", 0, 1, 'R')
    # Row 5
    pdf.cell(65, 5, "Finance, Closing & Buydown:", 0, 0, 'L')
    pdf.cell(30, 5, f"${btr_finance_closing:,.0f}", 0, 1, 'R')
    # Row 6
    pdf.cell(65, 5, "Built-in Developer Margin:", 0, 0, 'L')
    pdf.cell(30, 5, f"${developer_margin:,.0f}", 0, 1, 'R')
    # Row 7
    pdf.set_font("Arial", 'B', 9)
    pdf.cell(65, 5, "Total Appraised Value (ARV):", 0, 0, 'L')
    pdf.cell(30, 5, f"${total_arv:,.0f}", 0, 1, 'R')
    
    pdf.ln(5)
    
    # Transaction Table
    pdf.set_font("Arial", 'B', 9)
    pdf.cell(50, 6, "Phase / Cash Event", 1, 0, 'C')
    pdf.cell(110, 6, "Transaction Scope", 1, 0, 'C')
    pdf.cell(30, 6, "Outlay", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 8)
    transactions = [
        ("1. Horizontal Dev & Land", "Out-of-pocket acquisition and civil infrastructure", f"-${total_land_default:,.0f}"),
        ("2. Vertical Const. Cost", "Baseline Hard Costs + Indirects + GC Fee", f"-${total_const:,.0f}"),
        ("3. Taps & Soft Costs", "Water/sewer/electric/gas tie-ins, permits", f"-${total_vertical_soft:,.0f}"),
        ("4. Const. Finance Costs", f"Lender fees (${const_closing_fee:,.0f}) & interest (${carry_int_base:,.0f})", f"-${const_closing_fee + carry_int_base:,.0f}"),
        ("5. Perm. Takeout Fees", f"Comm. fees (${refi_closing_fee:,.0f}) + buydown (${default_buydown_cost:,.0f})", f"-${refi_closing_fee + default_buydown_cost:,.0f}")
    ]
    for e, s, a in transactions:
        pdf.cell(50, 6, e, 1, 0, 'L')
        pdf.cell(110, 6, s, 1, 0, 'L')
        pdf.cell(30, 6, a, 1, 1, 'R')
        
    pdf.set_font("Arial", 'B', 9)
    pdf.cell(160, 6, "Total Project Capital Basis:", 1, 0, 'L')
    pdf.cell(30, 6, f"-${total_project_basis:,.0f}", 1, 1, 'R')
    pdf.ln(5)

    # 7. OPERATING PERFORMANCE & DSCR
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 7. Operating Performance Summary", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    clean_op_summary = operating_summary_text.replace("**", "").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_op_summary)
    pdf.ln(3)

    pdf.set_font("Arial", 'B', 9)
    pdf.cell(80, 6, "Pro Forma Line Item", 1, 0, 'C')
    pdf.cell(55, 6, "Monthly", 1, 0, 'C')
    pdf.cell(55, 6, "Annual", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 8)
    for i in range(len(dscr_summary_data["Pro Forma Line Item"])):
        item = str(dscr_summary_data["Pro Forma Line Item"][i])
        mo = str(dscr_summary_data["Monthly"][i])
        yr = str(dscr_summary_data["Annual"][i])
        pdf.cell(80, 6, item, 1, 0, 'L')
        pdf.cell(55, 6, mo, 1, 0, 'R')
        pdf.cell(55, 6, yr, 1, 1, 'R')
    pdf.ln(5)

    # 8. REFINANCE CASH WATERFALL
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 8. Refinance Cash Waterfall (Payoff & Cash-Out)", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    clean_waterfall_summary = waterfall_summary_text.replace("**", "").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_waterfall_summary)
    pdf.ln(3)
    
    pdf.cell(130, 7, "(+) Permanent Takeout Loan Proceeds:", 0, 0)
    pdf.cell(60, 7, f"${loan_total:,.0f}", 0, 1, 'R')
    pdf.cell(130, 7, "(-) Construction Loan Payoff (Prin. & Int.):", 0, 0)
    pdf.cell(60, 7, f"-${actual_const_loan:,.0f}", 0, 1, 'R')
    pdf.cell(130, 7, "(-) Refinance Closing Fees & Buydown Points:", 0, 0)
    pdf.cell(60, 7, f"-${refi_closing_fee + default_buydown_cost:,.0f}", 0, 1, 'R')
    
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(130, 7, "= Net Cash Distributed to Sponsor at Closing:", 0, 0)
    pdf.cell(60, 7, f"${net_cash_at_closing:,.0f}", 0, 1, 'R')
    pdf.set_font("Arial", '', 10)
    pdf.cell(130, 7, "(-) Initial Sponsor Seed Capital Invested:", 0, 0)
    pdf.cell(60, 7, f"-${seed_capital:,.0f}", 0, 1, 'R')
    
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(130, 7, "= Final Tax-Free Cash Surplus / (Trapped Capital):", 0, 0)
    pdf.cell(60, 7, f"${cash_surplus:,.0f}", 0, 1, 'R')
    pdf.ln(5)

    # 9. STRATEGY COMPARISON: RETAIL VS BTR
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 9. Strategy Comparison: Retail Sell vs. Build-to-Rent", ln=1, fill=True)
    pdf.set_font("Arial", '', 9)
    
    # Text summary for PDF
    pdf.set_font("Arial", '', 10)
    clean_strategy_summary = strategy_comparison_text.replace("**", "").replace(r"\$", "$").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_strategy_summary)
    pdf.ln(3)

    # Table Header
    pdf.set_font("Arial", 'B', 9)
    pdf.cell(70, 6, "Financial Metric", 1, 0, 'C')
    pdf.cell(60, 6, "National Builder (Retail Sell)", 1, 0, 'C')
    pdf.cell(60, 6, "Build-to-Rent (DSCR Takeout)", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 8)
    for i in range(len(strategy_data["Financial Metric"])):
        metric = str(strategy_data["Financial Metric"][i])
        retail = str(strategy_data["National Builder (Retail Sell)"][i])
        btr = str(strategy_data["Build-to-Rent (Commercial DSCR Takeout)"][i])
        
        # Abbreviate for PDF constraints
        if "(~8% Realtor & concessions)" in retail:
            retail = retail.replace("(~8% Realtor & concessions)", "(8% Fees)")
        if "(Taxable ordinary income)" in retail:
            retail = retail.replace("(Taxable ordinary income)", "(Taxable)")
        if "(Const. finance, takeout, buydown)" in btr:
            btr = btr.replace("(Const. finance, takeout, buydown)", "(Financing Fees)")
        if "LTV DSCR Loan" in btr:
            btr = btr.replace(" LTV DSCR Loan", " LTV)")
        if "(Seed Capital Retained)" in btr:
            btr = btr.replace("(Seed Capital Retained)", "(Retained)")
        if "(Tax-Free Cash Out)" in btr:
            btr = btr.replace("(Tax-Free Cash Out)", "(Cash Out)")
            
        pdf.cell(70, 6, metric[:45], 1, 0, 'L')
        pdf.cell(60, 6, retail[:45], 1, 0, 'C')
        pdf.cell(60, 6, btr[:45], 1, 1, 'C')
    
    pdf.ln(5)
    pdf.set_font("Arial", 'I', 8)
    clean_footnote = strategy_footnote_text.replace("*Note: ", "Note: ").replace("*", "").replace(r"\$", "$").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 5, clean_footnote)
    pdf.ln(5)
    
    # 10. MULTI-UNIT PHASE & ANNUAL PROGRAM ANALYSIS
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 10. Multi-Unit Phase & Annual Program Analysis", ln=1, fill=True)
    
    pdf.set_font("Arial", '', 10)
    clean_scaling_intro = scaling_intro.replace("**", "").replace(r"\$", "$").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_scaling_intro)
    pdf.ln(3)

    pdf.set_font("Arial", 'B', 8)
    pdf.cell(50, 6, "Portfolio Metric", 1, 0, 'C')
    pdf.cell(40, 6, "1-House Baseline", 1, 0, 'C')
    pdf.cell(40, 6, "3-House Phase", 1, 0, 'C')
    pdf.cell(40, 6, "6-House Annual", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 7)
    for i in range(len(scaling_data["Portfolio Metric"])):
        metric = str(scaling_data["Portfolio Metric"][i])
        val1 = str(scaling_data["Single-House Baseline"][i])
        val3 = str(scaling_data["3-House Simultaneous Phase"][i])
        val6 = str(scaling_data["6-House Annual Build (Year 1)"][i])
        
        # Abbreviate for PDF
        metric = metric.replace("Permanent Commercial Loan Proceeds", "Takeout Loan Proceeds")
        metric = metric.replace("Total Retained Asset Equity", "Retained Asset Equity")
        
        pdf.cell(50, 6, metric[:40], 1, 0, 'L')
        pdf.cell(40, 6, val1[:30], 1, 0, 'C')
        pdf.cell(40, 6, val3[:30], 1, 0, 'C')
        pdf.cell(40, 6, val6[:30], 1, 1, 'C')
        
    pdf.ln(4)
    pdf.set_font("Arial", 'I', 9)
    clean_scaling_takeaway = scaling_takeaway.replace("**", "").replace(r"\$", "$").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 5, clean_scaling_takeaway)
    pdf.ln(5)

    # 11. ENTERPRISE S-CURVE DRAW SCHEDULE
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 11. Enterprise S-Curve Draw Schedule", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    clean_scurve_summary = scurve_summary.replace("**", "").replace(r"\$", "$").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_scurve_summary)
    pdf.ln(3)
    
    pdf.set_font("Arial", 'B', 9)
    pdf.cell(30, 6, "Month", 1, 0, 'C')
    pdf.cell(40, 6, "Cum. % Complete", 1, 0, 'C')
    pdf.cell(40, 6, "Monthly Draw", 1, 0, 'C')
    pdf.cell(40, 6, "Capitalized Int.", 1, 0, 'C')
    pdf.cell(40, 6, "Total Drawn Bal.", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 9)
    for i, row in df_schedule.iterrows():
        pdf.cell(30, 6, str(row["Month"]), 1, 0, 'C')
        pdf.cell(40, 6, str(row["Cum. % Complete"]), 1, 0, 'C')
        pdf.cell(40, 6, f"${row['Monthly Draw ($)']:,.0f}", 1, 0, 'R')
        pdf.cell(40, 6, f"${row['Capitalized Interest ($)']:,.0f}", 1, 0, 'R')
        pdf.cell(40, 6, f"${row['Total Drawn Balance ($)']:,.0f}", 1, 1, 'R')
    pdf.ln(5)
    
    # 12. SENSITIVITY ANALYSIS (STRESS TEST)
    pdf.set_font("Arial", 'B', 12)
    pdf.set_fill_color(220, 220, 220)
    pdf.cell(0, 8, " 12. Sensitivity Analysis (Stress Testing)", ln=1, fill=True)
    pdf.set_font("Arial", '', 10)
    
    clean_stress_summary = stress_test_summary.replace("**", "").replace(r"\$", "$").encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 6, clean_stress_summary)
    pdf.ln(3)
    
    pdf.set_font("Arial", 'B', 8)
    pdf.cell(60, 6, "Stress Scenario", 1, 0, 'C')
    pdf.cell(25, 6, "Final ARV", 1, 0, 'C')
    pdf.cell(20, 6, "Refi Rate", 1, 0, 'C')
    pdf.cell(25, 6, "DSCR Status", 1, 0, 'C')
    pdf.cell(35, 6, "Net Cash Surplus", 1, 0, 'C')
    pdf.cell(30, 6, "Day-1 Wealth", 1, 1, 'C')
    
    pdf.set_font("Arial", '', 8)
    for i, row in df_stress.iterrows():
        # Remove emojis for PDF rendering
        dscr_clean = str(row["DSCR Status"]).replace("🟢 ", "").replace("🔴 ", "")
        surplus_clean = str(row["Net Cash Surplus / (Trapped)"]).replace("🟢 ", "").replace("🔴 ", "")
        
        pdf.cell(60, 6, str(row["Stress Scenario"])[:35], 1, 0, 'L')
        pdf.cell(25, 6, str(row["Final ARV"]), 1, 0, 'C')
        pdf.cell(20, 6, str(row["Refi Rate"]), 1, 0, 'C')
        pdf.cell(25, 6, dscr_clean, 1, 0, 'C')
        pdf.cell(35, 6, surplus_clean, 1, 0, 'R')
        pdf.cell(30, 6, str(row["Total Day-1 Wealth"]), 1, 1, 'R')
    pdf.ln(5)

    # 13. REVERSE-ENGINEERING BREAKDOWN
    if cost_calc_mode in ["Reverse-Engineer from Appraisal", "Reverse-Engineer from Primary Comp"]:
        pdf.set_font("Arial", 'B', 12)
        pdf.set_fill_color(220, 220, 220)
        pdf.cell(0, 8, " 13. Retail Comp & Appraisal Reverse-Engineering", ln=1, fill=True)
        pdf.set_font("Arial", '', 10)
        
        # Add dynamic narrative to PDF
        rev_eng_summary_pdf = (
            f"Based on the comp's isolated rates, your specific {sqft} SF project has an estimated ARV "
            f"of ${reference_price:,.0f}. We are reverse-engineering this value to isolate the pure direct "
            f"hard costs from finished land, builder corporate margins, indirects, and soft costs. This ensures "
            f"the physical build budget perfectly aligns with proven market economics."
        )
        pdf.multi_cell(0, 6, rev_eng_summary_pdf.encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(3)
        
        pdf.cell(130, 7, "Project Target ARV (Derived from Comp):", 0, 0)
        pdf.cell(60, 7, f"${reference_price:,.0f}", 0, 1, 'R')
        pdf.cell(130, 7, "(-) Finished Lot Cost:", 0, 0)
        pdf.cell(60, 7, f"-${reference_price * lot_cost_pct:,.0f}", 0, 1, 'R')
        pdf.cell(130, 7, "(-) Builder Profit (Pre-tax):", 0, 0)
        pdf.cell(60, 7, f"-${reference_price * profit_margin_pct:,.0f}", 0, 1, 'R')
        pdf.cell(130, 7, "(-) Overhead & General Expenses:", 0, 0)
        pdf.cell(60, 7, f"-${reference_price * overhead_pct:,.0f}", 0, 1, 'R')
        pdf.cell(130, 7, "(-) Sales & Marketing:", 0, 0)
        pdf.cell(60, 7, f"-${reference_price * sales_pct:,.0f}", 0, 1, 'R')
        pdf.cell(130, 7, "(-) Financing Costs:", 0, 0)
        pdf.cell(60, 7, f"-${reference_price * finance_pct:,.0f}", 0, 1, 'R')
        
        pdf.set_font("Arial", 'B', 10)
        pdf.cell(130, 7, "= Total Construction Budget (Direct + Indirect):", 0, 0)
        pdf.cell(60, 7, f"${total_construction_budget:,.0f}", 0, 1, 'R')
        pdf.set_font("Arial", '', 10)
        
        pdf.cell(130, 7, "(-) Indirects (Permits, Temp Site, GC Fee):", 0, 0)
        pdf.cell(60, 7, f"-${total_indirect_costs:,.0f}", 0, 1, 'R')
        
        pdf.set_font("Arial", 'B', 10)
        pdf.cell(130, 7, "= Direct Hard Cost Budget:", 0, 0)
        pdf.cell(60, 7, f"${target_direct_hard_cost:,.0f}", 0, 1, 'R')
        pdf.set_font("Arial", '', 10)
        
        pdf.cell(130, 7, "(-) Fixed Auxiliary & Elevation Costs:", 0, 0)
        pdf.cell(60, 7, f"-${our_aux_cost_total:,.0f}", 0, 1, 'R')
        
        pdf.set_font("Arial", 'B', 10)
        pdf.cell(130, 7, "= Available Budget for Heated Shell:", 0, 0)
        pdf.cell(60, 7, f"${target_heated_hard_cost:,.0f}", 0, 1, 'R')
        pdf.ln(5)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        pdf.output(tmp.name)
        with open(tmp.name, "rb") as f:
            return f.read()

# Render PDF Download Button
st.download_button(
    label="📄 Download Enterprise Report (PDF)",
    data=create_pdf(pdf_detail_mode),
    file_name=f"Wickboldt_Capital_ProForma_{report_date.replace(' ', '_').replace(',', '')}.pdf",
    mime="application/pdf",
    type="primary",
    use_container_width=True
)

# ==========================================
# --- 15. PRESENTATION DECK GENERATION ---
# ==========================================
st.markdown("### 📊 15. Export Presentation Deck (PPTX & PDF)")
st.info("Automatically generates a 4-slide executive summary deck suitable for investor and commercial lender reviews.")

def create_pptx():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Wickboldt Capital | BTR Pro Forma Analysis"
    subtitle.text = f"Project: {project_name if project_name else 'TBD'}\nAddress: {project_address if project_address else 'TBD'}\nPrepared: {report_date}"

    # 2. Exec Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = f"Project Scale: {units} Unit(s)"
    
    p = tf.add_paragraph()
    p.text = f"Total Appraised Value (ARV): ${total_arv:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Total Project Basis: ${total_project_basis:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Built-in Developer Margin: ${developer_margin:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Net Cash Surplus at Refi: ${cash_surplus:,.0f}" if cash_surplus >= 0 else f"Retained Seed Capital: ${-cash_surplus:,.0f}"

    # 3. Operating Metrics
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Operating & DSCR Metrics"
    tf = body_shape.text_frame
    tf.text = f"Gross Monthly Income: ${total_gross_monthly_income:,.0f}"
    
    p = tf.add_paragraph()
    p.text = f"Net Operating Income (NOI): ${annual_noi/12:,.0f} / mo"
    p = tf.add_paragraph()
    p.text = f"Permanent Debt Service (P&I): ${total_monthly_pi:,.0f} / mo"
    p = tf.add_paragraph()
    p.text = f"Monthly Net Cash Flow: ${monthly_cash_flow:,.0f} (${monthly_cash_flow_per_door:,.0f} per door)"
    p = tf.add_paragraph()
    p.text = f"Actual DSCR: {actual_dscr:.2f}x (Target: {target_dscr_rate:.2f}x)"

    # 4. Capital Ledger
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Capital Outlay Breakdown"
    tf = body_shape.text_frame
    tf.text = f"Horizontal Land Basis: ${total_land_default:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Vertical Direct Hard Costs: ${total_hard_cost:,.0f} (${blended_cost_per_sf:.2f} / SF)"
    p = tf.add_paragraph()
    p.text = f"Indirects & GC Fee: ${total_indirect_costs:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Soft Costs & Utilities: ${total_vertical_soft:,.0f}"
    p = tf.add_paragraph()
    p.text = f"Financing & Closing Costs: ${btr_finance_closing:,.0f}"

    # Save
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        with open(tmp.name, "rb") as f:
            return f.read()

def create_slide_pdf():
    # Landscape orientation, millimeters, Letter size
    pdf = FPDF(orientation='L', unit='mm', format='Letter')
    
    # 1. Title Slide
    pdf.add_page()
    pdf.set_font('Arial', 'B', 28)
    pdf.cell(0, 40, '', ln=1) # Vertical Spacer
    pdf.cell(0, 15, 'Wickboldt Capital | BTR Pro Forma Analysis', align='C', ln=1)
    pdf.set_font('Arial', 'I', 16)
    pdf.cell(0, 10, f"Project: {project_name if project_name else 'TBD'}", align='C', ln=1)
    pdf.cell(0, 10, f"Address: {project_address if project_address else 'TBD'}", align='C', ln=1)
    pdf.cell(0, 10, f"Prepared: {report_date}", align='C', ln=1)
    
    # 2. Exec Summary
    pdf.add_page()
    pdf.set_font('Arial', 'B', 24)
    pdf.cell(0, 20, 'Executive Summary', ln=1)
    pdf.line(10, 30, 270, 30) # Separator line
    pdf.set_font('Arial', '', 18)
    pdf.ln(10)
    pdf.cell(0, 12, f"- Project Scale: {units} Unit(s)", ln=1)
    pdf.cell(0, 12, f"- Total Appraised Value (ARV): ${total_arv:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Total Project Basis: ${total_project_basis:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Built-in Developer Margin: ${developer_margin:,.0f}", ln=1)
    cash_text = f"- Net Cash Surplus at Refi: ${cash_surplus:,.0f}" if cash_surplus >= 0 else f"- Retained Seed Capital: ${-cash_surplus:,.0f}"
    pdf.cell(0, 12, cash_text, ln=1)
    
    # 3. Operating Metrics
    pdf.add_page()
    pdf.set_font('Arial', 'B', 24)
    pdf.cell(0, 20, 'Operating & DSCR Metrics', ln=1)
    pdf.line(10, 30, 270, 30)
    pdf.set_font('Arial', '', 18)
    pdf.ln(10)
    pdf.cell(0, 12, f"- Gross Monthly Income: ${total_gross_monthly_income:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Net Operating Income (NOI): ${annual_noi/12:,.0f} / mo", ln=1)
    pdf.cell(0, 12, f"- Permanent Debt Service (P&I): ${total_monthly_pi:,.0f} / mo", ln=1)
    pdf.cell(0, 12, f"- Monthly Net Cash Flow: ${monthly_cash_flow:,.0f} (${monthly_cash_flow_per_door:,.0f} per door)", ln=1)
    pdf.cell(0, 12, f"- Actual DSCR: {actual_dscr:.2f}x (Target: {target_dscr_rate:.2f}x)", ln=1)
    
    # 4. Capital Ledger
    pdf.add_page()
    pdf.set_font('Arial', 'B', 24)
    pdf.cell(0, 20, 'Capital Outlay Breakdown', ln=1)
    pdf.line(10, 30, 270, 30)
    pdf.set_font('Arial', '', 18)
    pdf.ln(10)
    pdf.cell(0, 12, f"- Horizontal Land Basis: ${total_land_default:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Vertical Direct Hard Costs: ${total_hard_cost:,.0f} (${blended_cost_per_sf:.2f} / SF)", ln=1)
    pdf.cell(0, 12, f"- Indirects & GC Fee: ${total_indirect_costs:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Soft Costs & Utilities: ${total_vertical_soft:,.0f}", ln=1)
    pdf.cell(0, 12, f"- Financing & Closing Costs: ${btr_finance_closing:,.0f}", ln=1)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        pdf.output(tmp.name)
        with open(tmp.name, "rb") as f:
            return f.read()

col1, col2 = st.columns(2)

with col1:
    st.download_button(
        label="📊 Download Presentation (PowerPoint)",
        data=create_pptx(),
        file_name=f"Wickboldt_Capital_Deck_{report_date.replace(' ', '_').replace(',', '')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        type="secondary",
        use_container_width=True
    )

with col2:
    st.download_button(
        label="📄 Download Presentation (PDF Slides)",
        data=create_slide_pdf(),
        file_name=f"Wickboldt_Capital_Slide_Deck_{report_date.replace(' ', '_').replace(',', '')}.pdf",
        mime="application/pdf",
        type="secondary",
        use_container_width=True
    )
